#!/usr/bin/env python3
"""
CI smoke-runner for Riley AI.

- If your real pipeline modules are importable, it uses them.
- Otherwise, it uses safe fallbacks (no network calls) so CI still runs.
- Outputs a small JSONL with chunks + a manifest CSV.

Strictness:
- Set RILEY_CI_STRICT=1 to fail when fallbacks are used (recommended after wiring).
"""

from __future__ import annotations

import argparse
import csv
import json
import os
import re
import sys
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Dict, List, Tuple

# -----------------------
# Try to import your code
# -----------------------
HAVE_REAL_EXTRACT = HAVE_REAL_PRIVACY = HAVE_REAL_CHUNK = HAVE_REAL_CLASSIFY = False
STRICT = os.getenv("RILEY_CI_STRICT", "0") == "1"

def _try_imports():
    global HAVE_REAL_EXTRACT, HAVE_REAL_PRIVACY, HAVE_REAL_CHUNK, HAVE_REAL_CLASSIFY
    try:
        from pipeline.extract import read_pdf as _read_pdf  # type: ignore
        from pipeline.extract import read_docx as _read_docx  # type: ignore
        from pipeline.extract import read_pptx as _read_pptx  # type: ignore
        try:
            from pipeline.extract import read_txt as _read_txt  # type: ignore
        except Exception:
            _read_txt = None
        HAVE_REAL_EXTRACT = True
        if _read_txt is not None:
            globals().update(read_pdf=_read_pdf, read_docx=_read_docx, read_pptx=_read_pptx, read_txt=_read_txt)
        else:
            globals().update(read_pdf=_read_pdf, read_docx=_read_docx, read_pptx=_read_pptx)
    except Exception:
        pass

    try:
        from pipeline.privacy import detect_pii as _detect_pii  # type: ignore
        from pipeline.privacy import redact_text as _redact_text  # type: ignore
        HAVE_REAL_PRIVACY = True
        globals().update(detect_pii=_detect_pii, redact_text=_redact_text)
    except Exception:
        pass

    try:
        from pipeline.chunk import chunk_text as _chunk_text  # type: ignore
        HAVE_REAL_CHUNK = True
        globals().update(chunk_text=_chunk_text)
    except Exception:
        pass

    try:
        from pipeline.classify_llm import classify as _classify  # type: ignore
        HAVE_REAL_CLASSIFY = True
        globals().update(classify=_classify)
    except Exception:
        pass


_try_imports()

# -----------------------
# Fallback implementations
# -----------------------

def _fallback_read_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)

def _fallback_read_docx(path: Path) -> str:
    import docx
    d = docx.Document(str(path))
    return "\n".join(p.text for p in d.paragraphs)

def _fallback_read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def _fallback_read_txt(path: Path) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")

EMAIL_RE = re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b")
PHONE_RE = re.compile(r"(?<!\d)(?:\+?\d[\d\-\s]{7,}\d)")
ADDR_RE = re.compile(r"\b(\d{1,5}\s+[A-Za-z0-9.\s]+(?:Street|St|Avenue|Ave|Road|Rd|Boulevard|Blvd|Lane|Ln|Way))\b", re.IGNORECASE)

def _fallback_detect_pii(text: str) -> Dict[str, int]:
    return {
        "emails": len(EMAIL_RE.findall(text)),
        "phones": len(PHONE_RE.findall(text)),
        "addresses": len(ADDR_RE.findall(text)),
    }

def _fallback_redact_text(text: str) -> str:
    t = EMAIL_RE.sub("[EMAIL]", text)
    t = PHONE_RE.sub("[PHONE]", t)
    t = ADDR_RE.sub("[ADDR]", t)
    return t

def _approx_tokens(text: str) -> int:
    # rough ~4 chars per token
    n = max(1, len(text))
    return max(1, n // 4)

def _fallback_chunk_text(text: str, max_tokens: int = 400, overlap_tokens: int = 40) -> List[str]:
    approx = max_tokens * 4
    overlap_chars = overlap_tokens * 4
    if len(text) <= approx:
        return [text]
    out = []
    i = 0
    while i < len(text):
        end = min(len(text), i + approx)
        out.append(text[i:end])
        if end == len(text):
            break
        i = end - overlap_chars
        if i < 0:
            i = 0
    return out

TAXONOMY = ["Fundraising", "Comms", "Field", "Policy", "General"]

def _fallback_classify(text: str) -> str:
    text_low = text.lower()
    if any(k in text_low for k in ("donate", "fundraising", "small-dollar", "contribute")):
        return "Fundraising"
    if any(k in text_low for k in ("press", "media", "announce", "statement", "release")):
        return "Comms"
    if any(k in text_low for k in ("volunteer", "door knock", "canvass", "phonebank", "rally")):
        return "Field"
    if any(k in text_low for k in ("policy", "plan", "healthcare", "education", "tax", "climate")):
        return "Policy"
    return "General"

# Select active implementations
read_pdf = globals().get("read_pdf", _fallback_read_pdf)
read_docx = globals().get("read_docx", _fallback_read_docx)
read_pptx = globals().get("read_pptx", _fallback_read_pptx)
read_txt = globals().get("read_txt", _fallback_read_txt)
detect_pii = globals().get("detect_pii", _fallback_detect_pii)
redact_text = globals().get("redact_text", _fallback_redact_text)
chunk_text = globals().get("chunk_text", _fallback_chunk_text)

MOCK_LLM = os.getenv("MOCK_LLM", "1") == "1"
classify = globals().get("classify", _fallback_classify)

if STRICT and not (HAVE_REAL_EXTRACT and HAVE_REAL_PRIVACY and HAVE_REAL_CHUNK and HAVE_REAL_CLASSIFY):
    missing = []
    if not HAVE_REAL_EXTRACT: missing.append("pipeline.extract")
    if not HAVE_REAL_PRIVACY: missing.append("pipeline.privacy")
    if not HAVE_REAL_CHUNK: missing.append("pipeline.chunk")
    if not HAVE_REAL_CLASSIFY: missing.append("pipeline.classify_llm")
    sys.stderr.write(f"[RILEY_CI_STRICT=1] Missing real modules: {', '.join(missing)}\n")
    sys.exit(2)

# -----------------------
# Driver
# -----------------------

from dataclasses import dataclass

@dataclass
class ChunkRecord:
    doc_id: str
    source_path: str
    chunk_index: int
    text: str
    approx_tokens: int
    taxonomy_label: str
    pii: Dict[str, int]

def extract_any(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return read_pdf(path)
    if ext == ".docx":
        return read_docx(path)
    if ext in (".pptx", ".ppt"):
        return read_pptx(path)
    if ext in (".txt", ".md"):
        return read_txt(path)
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""

def run(input_dir: Path, out_dir: Path, max_tokens: int = 400, overlap_tokens: int = 40):
    out_dir.mkdir(parents=True, exist_ok=True)
    artifacts_dir = out_dir / "artifacts"
    artifacts_dir.mkdir(parents=True, exist_ok=True)

    chunks_path = out_dir / "chunks.jsonl"
    manifest_path = artifacts_dir / "manifest.csv"

    with open(chunks_path, "w", encoding="utf-8") as jf, open(manifest_path, "w", newline="", encoding="utf-8") as mf:
        writer = csv.DictWriter(mf, fieldnames=["doc_id", "source_path", "chars", "approx_tokens", "chunks", "emails", "phones", "addresses"])
        writer.writeheader()

        for p in sorted(input_dir.rglob("*")):
            if not p.is_file():
                continue

            text = extract_any(p)
            if not text.strip():
                continue

            pii = detect_pii(text)
            redacted = redact_text(text)
            pieces = chunk_text(redacted, max_tokens=max_tokens, overlap_tokens=overlap_tokens)

            doc_id = p.stem
            for idx, piece in enumerate(pieces):
                label = classify(piece)
                rec = ChunkRecord(
                    doc_id=doc_id,
                    source_path=str(p.relative_to(input_dir)),
                    chunk_index=idx,
                    text=piece,
                    approx_tokens=_approx_tokens(piece),
                    taxonomy_label=label,
                    pii=pii,
                )
                jf.write(json.dumps(asdict(rec), ensure_ascii=False) + "\n")

            writer.writerow({
                "doc_id": doc_id,
                "source_path": str(p.relative_to(input_dir)),
                "chars": len(text),
                "approx_tokens": _approx_tokens(text),
                "chunks": len(pieces),
                "emails": pii.get("emails", 0),
                "phones": pii.get("phones", 0),
                "addresses": pii.get("addresses", 0),
            })

    return chunks_path, manifest_path

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", type=Path, required=True, help="Folder with campaign files")
    ap.add_argument("--out", type=Path, required=True, help="Output folder")
    ap.add_argument("--max_tokens", type=int, default=400)
    ap.add_argument("--overlap_tokens", type=int, default=40)
    args = ap.parse_args()

    chunks_path, manifest_path = run(args.input, args.out, args.max_tokens, args.overlap_tokens)
    print(f"Chunks: {chunks_path}")
    print(f"Manifest: {manifest_path}")

if __name__ == "__main__":
    main()
