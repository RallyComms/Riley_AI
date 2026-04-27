import asyncio
import hashlib
import io
import logging
import csv
import json
import re
import uuid
import time
from datetime import datetime, timedelta, timezone
from typing import Any, Dict, List, Optional, Tuple

from fastapi import UploadFile, HTTPException
from fastapi.concurrency import run_in_threadpool
from qdrant_client import models as qdrant_models
from qdrant_client.http.models import PointStruct, Filter, FieldCondition, MatchValue
from google.cloud import tasks_v2
from google.api_core.exceptions import AlreadyExists

from app.core.config import get_settings
from app.services.genai_client import get_genai_client
from app.services.llm_cost_guardrail import enforce_monthly_llm_cost_guardrail
from app.services.ocr import (
    gcs_uri_from_url,
    is_image_ext,
    run_ocr,
    run_ocr_document_from_gcs,
)
from app.services.preview import is_office_or_html, generate_pdf_preview
from app.services.qdrant import vector_service
from app.services.storage import StorageService
from app.services.token_utils import estimate_tokens
from app.services.visual_understanding import summarize_visual_content

# Configure logging
logger = logging.getLogger(__name__)

# Constants for text extraction limits
MAX_CHARS_TOTAL = 20_000
CONTENT_PREVIEW_LENGTH = 1_000
XLSX_MAX_SHEETS = 3
XLSX_MAX_ROWS = 100
XLSX_MAX_COLS = 25
PPTX_MAX_SLIDES = 50
CSV_MAX_ROWS = 300
CSV_MAX_COLS = 40
CHUNK_SIZE_TOKENS = 850
CHUNK_OVERLAP_TOKENS = 120
MICRO_CHUNK_SIZE_TOKENS = 240
MICRO_CHUNK_OVERLAP_TOKENS = 48
MACRO_CHUNK_SIZE_TOKENS = 1100
MACRO_CHUNK_OVERLAP_TOKENS = 180
MIN_EXTRACTED_CHARS = 120
EMBEDDING_COST_PER_1K_TOKENS_USD = 0.0001
BM25_MODEL = "qdrant/bm25"
OCR_APPEND_MAX_CHARS = 2500
VISION_APPEND_MAX_CHARS = 600
MERGED_CHUNK_MAX_CHARS = 6000
MIN_MEANINGFUL_NATIVE_CHARS = 80
MIN_MEANINGFUL_TOTAL_CHARS = 120
MIN_PARTIAL_TOTAL_CHARS = 35
MAX_RETRYABLE_INGESTION_ATTEMPTS = 5

TEXT_NATIVE_EXTENSIONS = {
    "pdf", "docx", "doc", "pptx", "ppt", "txt", "md", "rtf",
    "html", "htm", "csv", "xlsx", "xls", "json", "tsv",
}

IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "webp", "tiff"}


class IngestionPermanentFailure(RuntimeError):
    """Raised when ingestion should stop retrying for this job."""

    def __init__(self, reason: str) -> None:
        super().__init__(reason)
        self.reason = reason


def _format_file_size(size_bytes: int) -> str:
    """Format file size in bytes to human-readable string.
    
    Args:
        size_bytes: File size in bytes
        
    Returns:
        Formatted string like "5.2 MB" or "1.5 KB"
    """
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    else:
        return f"{size_bytes / (1024 * 1024):.1f} MB"




async def extract_text(
    file_bytes: bytes, filename: str
) -> str:
    """Router function to extract text from various file types.
    
    Note: OCR is NOT performed during extraction. OCR must be explicitly
    requested via the dedicated POST /files/{file_id}/ocr endpoint.
    
    Args:
        file_bytes: File content as bytes
        filename: Original filename (used to determine file type)
        
    Returns:
        Extracted text string, truncated to MAX_CHARS_TOTAL
        Returns "[Binary file — no text extracted]" on failure
    """
    # Normalize file extension
    file_ext = filename.split(".")[-1].lower() if "." in filename else ""
    
    try:
        if file_ext == "pdf":
            return await _extract_text_from_pdf(file_bytes)
        elif file_ext == "docx":
            return await _extract_text_from_docx(file_bytes)
        elif file_ext == "doc":
            return await _extract_text_from_doc(file_bytes)
        elif file_ext == "pptx":
            return await _extract_text_from_pptx(file_bytes)
        elif file_ext == "ppt":
            return await _extract_text_from_ppt(file_bytes)
        elif file_ext == "xlsx":
            return await _extract_text_from_xlsx(file_bytes)
        elif file_ext == "xls":
            return await _extract_text_from_xls(file_bytes)
        elif file_ext in ("html", "htm"):
            return await _extract_text_from_html(file_bytes)
        elif file_ext in ("txt", "md"):
            return await _extract_text_from_plain_text(file_bytes)
        elif file_ext == "csv":
            return await _extract_text_from_csv_like(file_bytes, delimiter=",")
        elif file_ext == "tsv":
            return await _extract_text_from_csv_like(file_bytes, delimiter="\t")
        elif file_ext == "json":
            return await _extract_text_from_json(file_bytes)
        elif file_ext == "rtf":
            return await _extract_text_from_rtf(file_bytes)
        elif file_ext in ("png", "jpg", "jpeg", "webp", "tiff"):
            # OCR is never performed during extraction - must use dedicated endpoint
            return await _extract_text_from_image(file_bytes, enable_ocr=False)
        else:
            return "[Unsupported file type for text extraction]"
    except Exception as exc:
        logger.warning(f"Text extraction failed for {filename}: {type(exc).__name__}: {exc}")
        return "[Binary file — no text extracted]"


async def _extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from a PDF file using pypdf from file content in memory."""
    try:
        from pypdf import PdfReader
        
        def _read_pdf():
            pdf_file = io.BytesIO(file_content)
            reader = PdfReader(pdf_file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            # Cap at MAX_CHARS_TOTAL
            if len(text) > MAX_CHARS_TOTAL:
                text = text[:MAX_CHARS_TOTAL]
            return text
        
        # Run blocking I/O in threadpool
        return await run_in_threadpool(_read_pdf)
    except ImportError:
        logger.error("pypdf is not installed. Install with: pip install pypdf")
        return "[PDF extraction unavailable — pypdf not installed]"
    except Exception as exc:
        logger.warning(f"PDF extraction failed: {exc}")
        return "[Binary file — no text extracted]"


async def _extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from a DOCX file using python-docx from file content in memory.
    
    Args:
        file_content: The DOCX file content as bytes
        
    Returns:
        Extracted text, capped at 20,000 characters
        
    Raises:
        HTTPException: If extraction fails or library is not installed
    """
    try:
        from docx import Document
        
        def _read_docx():
            docx_file = io.BytesIO(file_content)
            doc = Document(docx_file)
            text_parts = []
            
            # Extract text from all paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text.strip())
            
            # Join paragraphs with newlines
            text = "\n".join(text_parts)
            
            # Cap at MAX_CHARS_TOTAL
            if len(text) > MAX_CHARS_TOTAL:
                text = text[:MAX_CHARS_TOTAL]
            
            return text
        
        # Run blocking I/O in threadpool
        return await run_in_threadpool(_read_docx)
    except ImportError:
        logger.error("python-docx is not installed. Install with: pip install python-docx")
        return "[DOCX extraction unavailable — python-docx not installed]"
    except Exception as exc:
        logger.warning(f"DOCX extraction failed: {exc}")
        return "[Binary file — no text extracted]"


async def _extract_text_from_xlsx(file_content: bytes) -> str:
    """Extract text from an XLSX file using openpyxl from file content in memory.
    
    Processes up to 3 sheets, 100 rows per sheet, and 25 columns per row.
    Values are joined as CSV-style text.
    
    Args:
        file_content: The XLSX file content as bytes
        
    Returns:
        Extracted text, capped at 20,000 characters
        
    Raises:
        HTTPException: If extraction fails or library is not installed
    """
    try:
        from openpyxl import load_workbook
        
        def _read_xlsx():
            xlsx_file = io.BytesIO(file_content)
            workbook = load_workbook(xlsx_file, data_only=True)
            text_parts = []
            
            # Process up to XLSX_MAX_SHEETS sheets
            sheet_count = 0
            for sheet_name in workbook.sheetnames:
                if sheet_count >= XLSX_MAX_SHEETS:
                    break
                
                sheet = workbook[sheet_name]
                sheet_text_parts = []
                
                # Process up to XLSX_MAX_ROWS rows
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                    if row_idx >= XLSX_MAX_ROWS:
                        break
                    
                    # Process up to XLSX_MAX_COLS columns per row
                    row_values = []
                    for col_idx, cell_value in enumerate(row):
                        if col_idx >= XLSX_MAX_COLS:
                            break
                        
                        # Convert cell value to string, handling None
                        if cell_value is not None:
                            row_values.append(str(cell_value))
                        else:
                            row_values.append("")
                    
                    # Join row values as CSV-style (comma-separated)
                    if any(row_values):  # Only add non-empty rows
                        sheet_text_parts.append(",".join(row_values))
                
                # Add sheet content with sheet name header
                if sheet_text_parts:
                    text_parts.append(f"Sheet: {sheet_name}")
                    text_parts.extend(sheet_text_parts)
                
                sheet_count += 1
            
            # Join all text parts with newlines
            text = "\n".join(text_parts)
            
            # Cap at MAX_CHARS_TOTAL
            if len(text) > MAX_CHARS_TOTAL:
                text = text[:MAX_CHARS_TOTAL]
            
            return text
        
        # Run blocking I/O in threadpool
        return await run_in_threadpool(_read_xlsx)
    except ImportError:
        logger.error("openpyxl is not installed. Install with: pip install openpyxl")
        return "[XLSX extraction unavailable — openpyxl not installed]"
    except Exception as exc:
        logger.warning(f"XLSX extraction failed: {exc}")
        return "[Binary file — no text extracted]"


async def _extract_text_from_doc(file_content: bytes) -> str:
    """Extract text from a DOC file using textract (legacy binary format).
    
    Args:
        file_content: The DOC file content as bytes
        
    Returns:
        Extracted text, capped at MAX_CHARS_TOTAL
    """
    try:
        import textract
        
        def _read_doc():
            doc_file = io.BytesIO(file_content)
            # textract can work with BytesIO
            text = textract.process(doc_file, extension='doc', encoding='utf-8')
            if isinstance(text, bytes):
                text = text.decode('utf-8', errors='ignore')
            # Cap at MAX_CHARS_TOTAL
            if len(text) > MAX_CHARS_TOTAL:
                text = text[:MAX_CHARS_TOTAL]
            return text
        
        # Run blocking I/O in threadpool
        return await run_in_threadpool(_read_doc)
    except ImportError:
        logger.error("textract is not installed. Install with: pip install textract")
        return "[DOC extraction unavailable — textract not installed]"
    except Exception as exc:
        logger.warning(f"DOC extraction failed: {exc}")
        return "[Binary file — no text extracted]"


async def _extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from a PPTX file using python-pptx.
    
    Processes up to PPTX_MAX_SLIDES slides, extracting text from shapes.
    
    Args:
        file_content: The PPTX file content as bytes
        
    Returns:
        Extracted text, capped at MAX_CHARS_TOTAL
    """
    try:
        from pptx import Presentation
        
        def _read_pptx():
            pptx_file = io.BytesIO(file_content)
            prs = Presentation(pptx_file)
            text_parts = []
            
            # Process up to PPTX_MAX_SLIDES slides
            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx >= PPTX_MAX_SLIDES:
                    break
                
                slide_text_parts = []
                slide_title = ""
                try:
                    title_shape = slide.shapes.title
                    if title_shape and getattr(title_shape, "text", None):
                        slide_title = str(title_shape.text).strip()
                except Exception:
                    slide_title = ""
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        shape_text = str(shape.text).strip()
                        if shape_text:
                            slide_text_parts.append(shape_text)

                notes_text = ""
                try:
                    if bool(getattr(slide, "has_notes_slide", False)):
                        notes_slide = slide.notes_slide
                        notes_frame = getattr(notes_slide, "notes_text_frame", None)
                        if notes_frame and getattr(notes_frame, "text", None):
                            notes_text = str(notes_frame.text).strip()
                except Exception:
                    notes_text = ""

                if slide_title:
                    slide_text_parts = [part for part in slide_text_parts if part != slide_title]
                    slide_text_parts.insert(0, f"Title: {slide_title}")
                if notes_text:
                    slide_text_parts.append(f"Speaker notes: {notes_text}")
                
                # Add slide content with slide number header
                if slide_text_parts:
                    text_parts.append(f"Slide {slide_idx + 1}:")
                    text_parts.extend(slide_text_parts)
            
            # Join all text parts with newlines
            text = "\n".join(text_parts)
            
            # Cap at MAX_CHARS_TOTAL
            if len(text) > MAX_CHARS_TOTAL:
                text = text[:MAX_CHARS_TOTAL]
            
            return text
        
        # Run blocking I/O in threadpool
        return await run_in_threadpool(_read_pptx)
    except ImportError:
        logger.error("python-pptx is not installed. Install with: pip install python-pptx")
        return "[PPTX extraction unavailable — python-pptx not installed]"
    except Exception as exc:
        logger.warning(f"PPTX extraction failed: {exc}")
        return "[Binary file — no text extracted]"


async def _extract_text_from_html(file_content: bytes) -> str:
    """Extract text from an HTML file using BeautifulSoup.
    
    Strips scripts, styles, and other non-content elements.
    
    Args:
        file_content: The HTML file content as bytes
        
    Returns:
        Extracted text, capped at MAX_CHARS_TOTAL
    """
    try:
        from bs4 import BeautifulSoup
        
        def _read_html():
            html_file = io.BytesIO(file_content)
            soup = BeautifulSoup(html_file, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text and clean up whitespace
            text = soup.get_text()
            # Clean up multiple newlines and spaces
            lines = [line.strip() for line in text.splitlines()]
            lines = [line for line in lines if line]  # Remove empty lines
            text = "\n".join(lines)
            
            # Cap at MAX_CHARS_TOTAL
            if len(text) > MAX_CHARS_TOTAL:
                text = text[:MAX_CHARS_TOTAL]
            
            return text
        
        # Run blocking I/O in threadpool
        return await run_in_threadpool(_read_html)
    except ImportError:
        logger.error("beautifulsoup4 is not installed. Install with: pip install beautifulsoup4")
        return "[HTML extraction unavailable — beautifulsoup4 not installed]"
    except Exception as exc:
        logger.warning(f"HTML extraction failed: {exc}")
        return "[Binary file — no text extracted]"


async def _extract_text_from_plain_text(file_content: bytes) -> str:
    """Extract text from plain UTF-family text files."""
    def _read_text() -> str:
        text = file_content.decode("utf-8", errors="ignore")
        return text[:MAX_CHARS_TOTAL]
    return await run_in_threadpool(_read_text)


async def _extract_text_from_csv_like(file_content: bytes, delimiter: str) -> str:
    """Extract text from CSV/TSV files as row-wise comma-separated lines."""
    def _read_csv() -> str:
        decoded = file_content.decode("utf-8", errors="ignore")
        stream = io.StringIO(decoded)
        reader = csv.reader(stream, delimiter=delimiter)
        lines: List[str] = []
        for row_idx, row in enumerate(reader):
            if row_idx >= CSV_MAX_ROWS:
                break
            clipped = [str(value) for value in row[:CSV_MAX_COLS]]
            if any(cell.strip() for cell in clipped):
                lines.append(", ".join(clipped))
        text = "\n".join(lines)
        return text[:MAX_CHARS_TOTAL]
    return await run_in_threadpool(_read_csv)


async def _extract_text_from_json(file_content: bytes) -> str:
    """Extract text from JSON by pretty-printing parsed content."""
    def _read_json() -> str:
        decoded = file_content.decode("utf-8", errors="ignore")
        try:
            obj = json.loads(decoded)
            normalized = json.dumps(obj, ensure_ascii=True, indent=2)
        except Exception:
            normalized = decoded
        return normalized[:MAX_CHARS_TOTAL]
    return await run_in_threadpool(_read_json)


async def _extract_text_from_rtf(file_content: bytes) -> str:
    """Extract text from RTF using striprtf if available, fallback regex cleanup."""
    def _read_rtf() -> str:
        raw = file_content.decode("utf-8", errors="ignore")
        try:
            from striprtf.striprtf import rtf_to_text
            cleaned = rtf_to_text(raw)
        except Exception:
            cleaned = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", raw)
            cleaned = cleaned.replace("{", " ").replace("}", " ")
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        return cleaned[:MAX_CHARS_TOTAL]
    return await run_in_threadpool(_read_rtf)


async def _extract_text_from_ppt(file_content: bytes) -> str:
    """Extract text from legacy PPT via textract."""
    try:
        import textract

        def _read_ppt() -> str:
            payload = textract.process(io.BytesIO(file_content), extension="ppt", encoding="utf-8")
            if isinstance(payload, bytes):
                payload = payload.decode("utf-8", errors="ignore")
            return payload[:MAX_CHARS_TOTAL]

        return await run_in_threadpool(_read_ppt)
    except ImportError:
        return "[PPT extraction unavailable — textract not installed]"
    except Exception as exc:
        logger.warning(f"PPT extraction failed: {exc}")
        return "[Binary file — no text extracted]"


async def _extract_text_from_xls(file_content: bytes) -> str:
    """Extract text from legacy XLS via textract."""
    try:
        import textract

        def _read_xls() -> str:
            payload = textract.process(io.BytesIO(file_content), extension="xls", encoding="utf-8")
            if isinstance(payload, bytes):
                payload = payload.decode("utf-8", errors="ignore")
            return payload[:MAX_CHARS_TOTAL]

        return await run_in_threadpool(_read_xls)
    except ImportError:
        return "[XLS extraction unavailable — textract not installed]"
    except Exception as exc:
        logger.warning(f"XLS extraction failed: {exc}")
        return "[Binary file — no text extracted]"


async def _extract_text_from_image(file_content: bytes, enable_ocr: bool = False) -> str:
    """Extract text from an image file using OCR (if enabled).
    
    NOTE: This function is deprecated in favor of the OCR service.
    For images, we now return placeholder text and handle OCR separately.
    
    Args:
        file_content: The image file content as bytes
        enable_ocr: If True, perform OCR using pytesseract (deprecated)
        
    Returns:
        Placeholder text - OCR is handled separately via OCR endpoint
    """
    # Images should not have OCR run during upload
    # OCR is handled via dedicated endpoint
    return "[Image file — OCR not requested]"


async def _generate_embedding(text: str) -> List[float]:
    """Generate an embedding vector for the given text using Google Gemini.
    
    This function is used for:
    - Initial file upload embeddings
    - OCR re-embedding after OCR completes (in files.py OCR endpoint)
    
    Uses the new google-genai SDK via the shared client.
    """
    settings = get_settings()
    model_name = settings.EMBEDDING_MODEL
    
    def _embed_sync(truncated_text: str) -> List[float]:
        """Synchronous embedding function to run in threadpool."""
        try:
            # Get client (lazy initialization)
            client = get_genai_client()
            # Use the shared client to generate embedding
            response = client.models.embed_content(
                model=model_name,
                contents=truncated_text
            )
            
            # Extract embedding vector from response
            if not response.embeddings or len(response.embeddings) == 0:
                raise RuntimeError(
                    f"Embedding response did not contain any embeddings for model '{model_name}'."
                )
            
            embedding = response.embeddings[0].values
            if not isinstance(embedding, list):
                raise RuntimeError(
                    f"Embedding response did not contain a valid vector for model '{model_name}'."
                )
            
            return embedding
        except RuntimeError:
            # Re-raise RuntimeError as-is
            raise
        except Exception as exc:
            error_msg = (
                f"Embedding failed using model '{model_name}': {exc}"
            )
            logger.error(error_msg)
            raise RuntimeError(error_msg) from exc
    
    try:
        await enforce_monthly_llm_cost_guardrail()
        # Truncate text to 9000 characters to avoid token limits
        truncated_text = text[:9000]
        
        # Run blocking GenAI call in threadpool to avoid blocking event loop
        return await run_in_threadpool(_embed_sync, truncated_text)
    except RuntimeError:
        # Re-raise RuntimeError as-is
        raise
    except Exception as exc:
        error_msg = (
            f"Embedding failed using model '{model_name}': {exc}"
        )
        logger.error(error_msg)
        raise HTTPException(
            status_code=500,
            detail=error_msg
        ) from exc


def _estimate_tokens(text: str) -> int:
    """Rough token estimate used for chunk/cost telemetry."""
    return estimate_tokens(text)


def _chunk_point_uuid(parent_file_id: str, chunk_type: str, chunk_index: int) -> str:
    """Deterministic UUID for chunk point IDs (Qdrant-safe)."""
    return str(uuid.uuid5(uuid.NAMESPACE_URL, f"riley-chunk:{parent_file_id}:{chunk_type}:{chunk_index}"))


def _clean_extracted_text(text: str) -> str:
    """Normalize whitespace to improve quality checks and chunking stability."""
    normalized = (text or "").replace("\x00", " ").strip()
    normalized = re.sub(r"\r\n?", "\n", normalized)
    normalized = re.sub(r"[ \t]+", " ", normalized)
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    return normalized


def _assess_extraction_quality(text: str, filename: str, file_size_bytes: int) -> Tuple[str, str]:
    """Classify native extraction quality and return (status, reason)."""
    cleaned = _clean_extracted_text(text)
    ext = filename.split(".")[-1].lower() if "." in filename else ""

    if not cleaned:
        if ext in {"pdf", "pptx", "ppt"}:
            return "ocr_needed", "No extractable native text detected for visual document"
        return "low_text", "No extractable text detected"

    lowered = cleaned.lower()
    placeholder_markers = [
        "[binary file",
        "no text extracted",
        "extraction unavailable",
        "unsupported file type",
    ]
    if any(marker in lowered for marker in placeholder_markers):
        return "low_text", "Extraction returned placeholder text"

    # Keep native gate permissive; final status now uses multimodal scoring.
    if len(cleaned) < MIN_MEANINGFUL_NATIVE_CHARS and file_size_bytes > 10_000:
        if ext in {"pdf", "pptx", "ppt"}:
            return "ocr_needed", "Native text is weak for visual-heavy document; OCR recommended"
        return "low_text", "Extracted text too short for file size"

    if len(cleaned) < MIN_PARTIAL_TOTAL_CHARS:
        if ext in {"pdf", "pptx", "ppt"}:
            return "ocr_needed", "Native text is very short; OCR recommended"
        return "low_text", "Native extraction too short"

    return "indexed", "Native extraction sufficient"


def _compute_content_signals(
    *,
    native_text: str,
    merged_text: str,
    segments: List[Dict[str, Any]],
) -> Dict[str, int]:
    native_char_count = len(_clean_extracted_text(native_text))
    merged_char_count = len(_clean_extracted_text(merged_text))
    ocr_char_count = 0
    vision_caption_count = 0
    vision_caption_chars = 0
    all_slots: set[str] = set()
    nonempty_slots: set[str] = set()

    for idx, segment in enumerate(segments):
        text = _clean_extracted_text(str(segment.get("text") or segment.get("raw_text") or ""))
        has_ocr = bool(segment.get("ocr_text_present")) or str(segment.get("ocr_status") or "").lower() == "complete"
        if has_ocr and text:
            ocr_char_count += len(text)

        caption = str(segment.get("vision_caption") or "").strip()
        if caption:
            vision_caption_count += 1
            vision_caption_chars += len(caption)

        loc_type = str(segment.get("location_type") or "").lower()
        loc_value = str(segment.get("location_value") or "").strip()
        if loc_type in {"page", "slide", "image_page"} and loc_value:
            slot = f"{loc_type}:{loc_value}"
        else:
            slot = f"segment:{idx + 1}"
        all_slots.add(slot)
        if text or caption:
            nonempty_slots.add(slot)

    total_page_count = len(all_slots)
    nonempty_page_count = len(nonempty_slots)
    return {
        "native_char_count": native_char_count,
        "ocr_char_count": ocr_char_count,
        "vision_caption_count": vision_caption_count,
        "vision_caption_chars": vision_caption_chars,
        "merged_char_count": merged_char_count,
        "nonempty_page_count": nonempty_page_count,
        "total_page_count": total_page_count,
    }


def _decide_final_ingestion_status(
    *,
    file_type: str,
    quality_status: str,
    quality_reason: str,
    signals: Dict[str, int],
) -> Tuple[str, str, str]:
    native_chars = int(signals.get("native_char_count") or 0)
    ocr_chars = int(signals.get("ocr_char_count") or 0)
    merged_chars = int(signals.get("merged_char_count") or 0)
    nonempty_pages = int(signals.get("nonempty_page_count") or 0)
    total_pages = int(signals.get("total_page_count") or 0)
    vision_caption_count = int(signals.get("vision_caption_count") or 0)
    vision_caption_chars = int(signals.get("vision_caption_chars") or 0)

    if (
        merged_chars >= MIN_MEANINGFUL_TOTAL_CHARS
        or native_chars >= MIN_MEANINGFUL_TOTAL_CHARS
        or ocr_chars >= MIN_MEANINGFUL_TOTAL_CHARS
        or (nonempty_pages >= 2 and merged_chars >= 70)
    ):
        return "indexed", "indexed_multichannel_sufficient", "indexed: native/OCR/merged content is sufficiently usable"

    if (
        merged_chars >= MIN_PARTIAL_TOTAL_CHARS
        or native_chars >= MIN_PARTIAL_TOTAL_CHARS
        or ocr_chars >= MIN_PARTIAL_TOTAL_CHARS
        or (vision_caption_count >= 2 and vision_caption_chars >= 80)
        or (nonempty_pages >= 1 and total_pages >= 1 and (merged_chars + vision_caption_chars) >= 50)
    ):
        return "partial", "partial_multichannel_limited", "partial: meaningful content exists but coverage/quality is limited"

    if quality_status == "ocr_needed":
        return "low_text", "low_text_all_channels_weak", "low_text: native extraction weak and OCR/vision did not recover enough usable content"
    return "low_text", "low_text_all_channels_weak", "low_text: all content channels are insufficient"


def _chunk_text_by_tokens(text: str, chunk_size_tokens: int = CHUNK_SIZE_TOKENS, overlap_tokens: int = CHUNK_OVERLAP_TOKENS) -> List[str]:
    """Split text into overlapping token-ish chunks using whitespace tokenization."""
    normalized = _clean_extracted_text(text)
    tokens = normalized.split()
    if not tokens:
        return []

    chunks: List[str] = []
    step = max(1, chunk_size_tokens - overlap_tokens)
    start = 0
    while start < len(tokens):
        end = min(len(tokens), start + chunk_size_tokens)
        chunk = " ".join(tokens[start:end]).strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(tokens):
            break
        start += step
    return chunks


def _fallback_segments_from_text(text: str, file_type: str) -> List[Dict[str, Any]]:
    normalized = (text or "").strip()
    if not normalized:
        return []
    parts = [part.strip() for part in re.split(r"\n{2,}", normalized) if part.strip()]
    if not parts:
        parts = [normalized]
    segments: List[Dict[str, Any]] = []
    for idx, part in enumerate(parts):
        segments.append(
            {
                "text": part,
                "raw_text": part,
                "location_type": "paragraph",
                "location_value": str(idx + 1),
                "section_path": f"{file_type}:body",
            }
        )
    return segments


async def _extract_structured_segments(
    file_bytes: bytes,
    filename: str,
    extracted_text: str,
) -> List[Dict[str, Any]]:
    """Best-effort structural segmentation for citation-grade chunk metadata."""
    file_ext = filename.split(".")[-1].lower() if "." in filename else ""

    try:
        if file_ext == "pdf":
            from pypdf import PdfReader

            def _read_pdf_segments() -> List[Dict[str, Any]]:
                reader = PdfReader(io.BytesIO(file_bytes))
                result: List[Dict[str, Any]] = []
                for page_idx, page in enumerate(reader.pages):
                    text = (page.extract_text() or "").strip()
                    if not text:
                        continue
                    result.append(
                        {
                            "text": text,
                            "raw_text": text,
                            "location_type": "page",
                            "location_value": str(page_idx + 1),
                            "section_path": f"page:{page_idx + 1}",
                        }
                    )
                return result

            segments = await run_in_threadpool(_read_pdf_segments)
            if segments:
                return segments

        if file_ext == "pptx":
            from pptx import Presentation

            def _read_pptx_segments() -> List[Dict[str, Any]]:
                prs = Presentation(io.BytesIO(file_bytes))
                result: List[Dict[str, Any]] = []
                total_slides = min(len(prs.slides), PPTX_MAX_SLIDES)
                skipped_slides = 0
                for slide_idx, slide in enumerate(prs.slides):
                    if slide_idx >= PPTX_MAX_SLIDES:
                        break
                    parts: List[str] = []
                    slide_title = ""
                    try:
                        title_shape = slide.shapes.title
                        if title_shape and getattr(title_shape, "text", None):
                            slide_title = str(title_shape.text).strip()
                    except Exception:
                        slide_title = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text = shape.text.strip()
                            if text:
                                parts.append(text)
                    notes_text = ""
                    try:
                        if bool(getattr(slide, "has_notes_slide", False)):
                            notes_slide = slide.notes_slide
                            notes_frame = getattr(notes_slide, "notes_text_frame", None)
                            if notes_frame and getattr(notes_frame, "text", None):
                                notes_text = str(notes_frame.text).strip()
                    except Exception:
                        notes_text = ""
                    if slide_title:
                        parts = [part for part in parts if part != slide_title]
                        parts.insert(0, slide_title)
                    if notes_text:
                        parts.append(f"Speaker notes: {notes_text}")
                    if not parts:
                        skipped_slides += 1
                        continue
                    merged = "\n".join(parts)
                    result.append(
                        {
                            "text": merged,
                            "raw_text": merged,
                            "location_type": "slide",
                            "location_value": str(slide_idx + 1),
                            "section_path": f"slide:{slide_idx + 1}",
                            "slide_title": slide_title or None,
                            "slide_number": slide_idx + 1,
                            "pptx_total_slides": total_slides,
                            "pptx_skipped_slides": skipped_slides,
                        }
                    )
                if result:
                    final_skipped = max(0, total_slides - len(result))
                    for item in result:
                        item["pptx_total_slides"] = total_slides
                        item["pptx_skipped_slides"] = final_skipped
                return result

            segments = await run_in_threadpool(_read_pptx_segments)
            if segments:
                return segments

        if file_ext == "docx":
            from docx import Document

            def _read_docx_segments() -> List[Dict[str, Any]]:
                doc = Document(io.BytesIO(file_bytes))
                result: List[Dict[str, Any]] = []
                heading_stack: List[str] = []
                para_index = 0
                for paragraph in doc.paragraphs:
                    text = (paragraph.text or "").strip()
                    if not text:
                        continue
                    para_index += 1
                    style_name = (paragraph.style.name or "").strip().lower() if paragraph.style else ""
                    if style_name.startswith("heading"):
                        heading_stack.append(text)
                        heading_stack = heading_stack[-4:]
                    section_path = " > ".join(heading_stack) if heading_stack else "body"
                    result.append(
                        {
                            "text": text,
                            "raw_text": text,
                            "location_type": "paragraph",
                            "location_value": str(para_index),
                            "section_path": section_path,
                        }
                    )
                return result

            segments = await run_in_threadpool(_read_docx_segments)
            if segments:
                return segments

        if file_ext in {"xlsx"}:
            from openpyxl import load_workbook

            def _read_xlsx_segments() -> List[Dict[str, Any]]:
                wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
                result: List[Dict[str, Any]] = []
                for sheet_name in wb.sheetnames[:XLSX_MAX_SHEETS]:
                    sheet = wb[sheet_name]
                    for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                        if row_idx > XLSX_MAX_ROWS:
                            break
                        values = [str(v) for v in row[:XLSX_MAX_COLS] if v is not None and str(v).strip()]
                        if not values:
                            continue
                        text = ", ".join(values)
                        result.append(
                            {
                                "text": text,
                                "raw_text": text,
                                "location_type": "sheet_row",
                                "location_value": f"{sheet_name}:{row_idx}",
                                "section_path": sheet_name,
                            }
                        )
                return result

            segments = await run_in_threadpool(_read_xlsx_segments)
            if segments:
                return segments

        if file_ext in {"csv", "tsv"}:
            delimiter = "," if file_ext == "csv" else "\t"

            def _read_csv_segments() -> List[Dict[str, Any]]:
                decoded = file_bytes.decode("utf-8", errors="ignore")
                stream = io.StringIO(decoded)
                reader = csv.reader(stream, delimiter=delimiter)
                result: List[Dict[str, Any]] = []
                for row_idx, row in enumerate(reader, start=1):
                    if row_idx > CSV_MAX_ROWS:
                        break
                    values = [str(v) for v in row[:CSV_MAX_COLS] if str(v).strip()]
                    if not values:
                        continue
                    text = ", ".join(values)
                    result.append(
                        {
                            "text": text,
                            "raw_text": text,
                            "location_type": "table_row",
                            "location_value": str(row_idx),
                            "section_path": "table",
                        }
                    )
                return result

            segments = await run_in_threadpool(_read_csv_segments)
            if segments:
                return segments

    except Exception as exc:
        logger.warning("structured_segment_extraction_failed filename=%s error=%s", filename, exc)

    return _fallback_segments_from_text(extracted_text, file_ext or "unknown")


def _prepare_segments_for_chunking(segments: List[Dict[str, Any]]) -> Tuple[str, List[Dict[str, Any]]]:
    prepared: List[Dict[str, Any]] = []
    merged_parts: List[str] = []
    cursor = 0

    for segment in segments:
        raw_text = str(segment.get("raw_text") or segment.get("text") or "")
        cleaned = _clean_extracted_text(raw_text)
        if not cleaned:
            continue
        if merged_parts:
            merged_parts.append("\n\n")
            cursor += 2
        char_start = cursor
        merged_parts.append(cleaned)
        cursor += len(cleaned)
        char_end = cursor
        prepared.append(
            {
                "text": cleaned,
                "raw_text": raw_text,
                "location_type": str(segment.get("location_type") or "paragraph"),
                "location_value": str(segment.get("location_value") or ""),
                "section_path": str(segment.get("section_path") or ""),
                "slide_title": str(segment.get("slide_title") or "").strip() or None,
                "slide_number": (
                    int(segment.get("slide_number"))
                    if isinstance(segment.get("slide_number"), int)
                    else (
                        int(str(segment.get("location_value") or "").strip())
                        if str(segment.get("location_value") or "").strip().isdigit()
                        and str(segment.get("location_type") or "").strip().lower() == "slide"
                        else None
                    )
                ),
                "pptx_total_slides": (
                    int(str(segment.get("pptx_total_slides") or "").strip())
                    if str(segment.get("pptx_total_slides") or "").strip().isdigit()
                    else None
                ),
                "pptx_skipped_slides": (
                    int(str(segment.get("pptx_skipped_slides") or "").strip())
                    if str(segment.get("pptx_skipped_slides") or "").strip().isdigit()
                    else None
                ),
                "char_start": char_start,
                "char_end": char_end,
            }
        )

    return "".join(merged_parts), prepared


def _chunk_text_with_offsets(
    text: str,
    *,
    chunk_size_tokens: int,
    overlap_tokens: int,
) -> List[Dict[str, Any]]:
    matches = list(re.finditer(r"\S+", text))
    if not matches:
        return []
    chunks: List[Dict[str, Any]] = []
    step = max(1, chunk_size_tokens - overlap_tokens)
    start_idx = 0
    while start_idx < len(matches):
        end_idx = min(len(matches), start_idx + chunk_size_tokens)
        char_start = matches[start_idx].start()
        char_end = matches[end_idx - 1].end()
        chunk_text = text[char_start:char_end].strip()
        if chunk_text:
            chunks.append(
                {
                    "text": chunk_text,
                    "char_start": char_start,
                    "char_end": char_end,
                }
            )
        if end_idx >= len(matches):
            break
        start_idx += step
    return chunks


def _split_text_with_char_limit_offsets(text: str, *, char_limit: int) -> List[Dict[str, Any]]:
    if char_limit <= 0:
        return []
    matches = list(re.finditer(r"\S+", text))
    if not matches:
        return []
    chunks: List[Dict[str, Any]] = []
    start_idx = 0
    while start_idx < len(matches):
        char_start = matches[start_idx].start()
        end_idx = start_idx
        while end_idx + 1 < len(matches):
            next_end = matches[end_idx + 1].end()
            if next_end - char_start > char_limit:
                break
            end_idx += 1
        char_end = matches[end_idx].end()
        chunk_text = text[char_start:char_end].strip()
        if chunk_text:
            chunks.append(
                {
                    "text": chunk_text,
                    "char_start": char_start,
                    "char_end": char_end,
                }
            )
        start_idx = end_idx + 1
    return chunks


def _build_pptx_slide_micro_chunks(
    segments: List[Dict[str, Any]],
    *,
    chunk_size_tokens: int,
    overlap_tokens: int,
    max_chars_per_chunk: int,
) -> Tuple[List[Dict[str, Any]], Dict[str, int]]:
    """Build slide-preserving micro chunks for PPTX content."""
    explicit_slide_segments: List[Dict[str, Any]] = []
    total_slides_hint = 0
    for segment in segments:
        total_raw = str(segment.get("pptx_total_slides") or "").strip()
        if total_raw.isdigit():
            total_slides_hint = max(total_slides_hint, int(total_raw))
        loc_type = str(segment.get("location_type") or "").strip().lower()
        loc_value = str(segment.get("location_value") or "").strip()
        segment_text = str(segment.get("text") or "").strip()
        if loc_type == "slide" and loc_value.isdigit() and segment_text:
            explicit_slide_segments.append(
                {
                    "slide_number": int(loc_value),
                    "slide_title": str(segment.get("slide_title") or "").strip() or None,
                    "text": segment_text,
                    "char_start": int(segment.get("char_start") or 0),
                    "char_end": int(segment.get("char_end") or 0),
                }
            )

    if not explicit_slide_segments:
        # Fallback path for degraded segmentation: infer slide blocks from inline headers.
        for idx, segment in enumerate(segments, start=1):
            source_text = str(segment.get("text") or "").strip()
            if not source_text:
                continue
            segment_start = int(segment.get("char_start") or 0)
            slide_header_matches = list(re.finditer(r"(?im)^\s*slide\s+(\d+)\s*:\s*", source_text))
            if slide_header_matches:
                for match_idx, match in enumerate(slide_header_matches):
                    content_start = match.end()
                    content_end = (
                        slide_header_matches[match_idx + 1].start()
                        if match_idx + 1 < len(slide_header_matches)
                        else len(source_text)
                    )
                    content_text = source_text[content_start:content_end].strip()
                    if not content_text:
                        continue
                    explicit_slide_segments.append(
                        {
                            "slide_number": int(match.group(1)),
                            "slide_title": None,
                            "text": content_text,
                            "char_start": segment_start + content_start,
                            "char_end": segment_start + content_end,
                        }
                    )
            else:
                explicit_slide_segments.append(
                    {
                        "slide_number": idx,
                        "slide_title": None,
                        "text": source_text,
                        "char_start": segment_start,
                        "char_end": int(segment.get("char_end") or segment_start + len(source_text)),
                    }
                )

    slide_best: Dict[int, Dict[str, Any]] = {}
    for entry in explicit_slide_segments:
        slide_number = int(entry.get("slide_number") or 0)
        if slide_number <= 0:
            continue
        existing = slide_best.get(slide_number)
        candidate_text = str(entry.get("text") or "")
        if existing is None or len(candidate_text) > len(str(existing.get("text") or "")):
            merged = dict(entry)
            if existing and not merged.get("slide_title"):
                merged["slide_title"] = existing.get("slide_title")
            slide_best[slide_number] = merged
        elif existing and not existing.get("slide_title") and entry.get("slide_title"):
            existing["slide_title"] = entry.get("slide_title")

    chunks: List[Dict[str, Any]] = []
    for slide_number in sorted(slide_best.keys()):
        entry = slide_best[slide_number]
        slide_text = str(entry.get("text") or "").strip()
        if not slide_text:
            continue
        token_chunks = _chunk_text_with_offsets(
            slide_text,
            chunk_size_tokens=chunk_size_tokens,
            overlap_tokens=overlap_tokens,
        )
        if not token_chunks:
            token_chunks = [{"text": slide_text, "char_start": 0, "char_end": len(slide_text)}]
        slide_subchunk_index = 0
        for token_chunk in token_chunks:
            token_text = str(token_chunk.get("text") or "").strip()
            if not token_text:
                continue
            token_start = int(token_chunk.get("char_start") or 0)
            char_chunks = _split_text_with_char_limit_offsets(
                token_text,
                char_limit=max_chars_per_chunk,
            )
            if not char_chunks:
                char_chunks = [{"text": token_text, "char_start": 0, "char_end": len(token_text)}]
            for char_chunk in char_chunks:
                chunk_text = str(char_chunk.get("text") or "").strip()
                if not chunk_text:
                    continue
                local_start = token_start + int(char_chunk.get("char_start") or 0)
                local_end = token_start + int(char_chunk.get("char_end") or len(chunk_text))
                chunks.append(
                    {
                        "text": chunk_text,
                        "char_start": int(entry.get("char_start") or 0) + local_start,
                        "char_end": int(entry.get("char_start") or 0) + local_end,
                        "slide_number": slide_number,
                        "slide_title": entry.get("slide_title"),
                        "slide_subchunk_index": slide_subchunk_index,
                    }
                )
                slide_subchunk_index += 1

    indexed_slides = len(slide_best)
    inferred_total = max(slide_best.keys()) if slide_best else 0
    total_slides = max(total_slides_hint, inferred_total)
    stats = {
        "total_slides": total_slides,
        "indexed_slides": indexed_slides,
        "skipped_slides": max(0, total_slides - indexed_slides),
    }
    return chunks, stats


def _best_segment_for_span(span_start: int, span_end: int, segments: List[Dict[str, Any]]) -> Dict[str, Any]:
    best: Dict[str, Any] = {}
    best_overlap = -1
    for segment in segments:
        seg_start = int(segment.get("char_start") or 0)
        seg_end = int(segment.get("char_end") or 0)
        overlap = max(0, min(span_end, seg_end) - max(span_start, seg_start))
        if overlap > best_overlap:
            best_overlap = overlap
            best = segment
    return best


def _ocr_required_for_file(
    *,
    file_type: str,
    quality_status: str,
    cleaned_text: str,
    file_size: int,
) -> bool:
    """Decide whether OCR should run for multimodal ingestion."""
    ext = (file_type or "").lower()
    ocr_capable_doc_types = {"pdf", "pptx", "ppt"}
    if ext in IMAGE_EXTENSIONS:
        return True
    if ext in ocr_capable_doc_types and quality_status in {"ocr_needed", "low_text"}:
        # P2 rule: low-text visual documents should always get an OCR attempt.
        return True
    if quality_status == "ocr_needed":
        # Safety net for any future OCR-capable file types.
        return True
    # Optional OCR for weak text on visual-heavy document types.
    if ext in ocr_capable_doc_types and (
        len(cleaned_text) < max(600, MIN_EXTRACTED_CHARS * 3) or file_size > 250_000
    ):
        return True
    return False


def _build_ocr_segments(
    *,
    filename: str,
    file_type: str,
    ocr_result: Dict[str, Any],
) -> List[Dict[str, Any]]:
    ext = (file_type or "").lower()
    pages = ocr_result.get("pages") or []
    segments: List[Dict[str, Any]] = []
    if pages:
        for entry in pages:
            text = str(entry.get("text") or "").strip()
            if not text:
                continue
            page_num = entry.get("page") or len(segments) + 1
            location_type = "page" if ext == "pdf" else "slide" if ext in {"pptx", "ppt"} else "image_page"
            segments.append(
                {
                    "text": text,
                    "raw_text": text,
                    "location_type": location_type,
                    "location_value": str(page_num),
                    "section_path": f"ocr:{location_type}:{page_num}",
                    "ocr_text_present": True,
                    "ocr_confidence": entry.get("confidence"),
                    "ocr_status": "complete",
                    "multimodal_status": "ocr_enriched",
                }
            )
        return segments

    merged = str(ocr_result.get("text") or "").strip()
    if not merged:
        return segments
    segments.append(
        {
            "text": merged,
            "raw_text": merged,
            "location_type": "image" if ext in IMAGE_EXTENSIONS else "ocr_text",
            "location_value": "1",
            "section_path": f"ocr:{filename}",
            "ocr_text_present": True,
            "ocr_confidence": ocr_result.get("confidence"),
            "ocr_status": "complete",
            "multimodal_status": "ocr_enriched",
        }
    )
    return segments


async def _render_pdf_pages_for_vision(
    *,
    pdf_bytes: bytes,
    page_numbers: List[int],
    max_pages: int,
) -> Dict[int, bytes]:
    """Render selected PDF pages into PNG bytes for multimodal captioning."""
    try:
        import fitz  # type: ignore
    except ImportError:
        logger.warning("vision_render_skipped reason=PyMuPDFNotInstalled")
        return {}
    if not page_numbers:
        return {}
    unique_pages = []
    seen = set()
    for page in page_numbers:
        if page in seen or page <= 0:
            continue
        seen.add(page)
        unique_pages.append(page)
        if len(unique_pages) >= max_pages:
            break
    rendered: Dict[int, bytes] = {}

    def _render_sync() -> Dict[int, bytes]:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        out: Dict[int, bytes] = {}
        try:
            for page_num in unique_pages:
                if page_num - 1 < 0 or page_num - 1 >= len(doc):
                    continue
                page = doc[page_num - 1]
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                out[page_num] = pix.tobytes("png")
        finally:
            doc.close()
        return out

    try:
        rendered = await run_in_threadpool(_render_sync)
    except Exception as exc:
        logger.warning("vision_render_failed error=%s", type(exc).__name__)
    return rendered


def _vision_candidate_pages(
    *,
    file_type: str,
    quality_status: str,
    segments: List[Dict[str, Any]],
    ocr_enabled: bool,
    max_pages: int,
) -> List[int]:
    ext = (file_type or "").lower()
    if ext in IMAGE_EXTENSIONS:
        return [1]
    if ext not in {"pdf", "pptx", "ppt"}:
        return []
    if quality_status == "indexed" and not ocr_enabled and ext == "pdf":
        # For clean text PDFs, skip default visual pass.
        return []
    candidates: List[int] = []
    for segment in segments:
        loc_type = str(segment.get("location_type") or "").lower()
        if loc_type not in {"page", "slide", "image_page"}:
            continue
        loc_val = str(segment.get("location_value") or "").strip()
        if not loc_val.isdigit():
            continue
        num = int(loc_val)
        if num <= 0:
            continue
        candidates.append(num)
        if len(candidates) >= max_pages * 2:
            break
    deduped: List[int] = []
    seen = set()
    for page in candidates:
        if page in seen:
            continue
        seen.add(page)
        deduped.append(page)
        if len(deduped) >= max_pages:
            break
    return deduped


def _merge_chunk_text(
    *,
    native_text: str,
    ocr_text: Optional[str],
    vision_caption: Optional[str],
) -> str:
    def _norm(value: str) -> str:
        lowered = value.lower().strip()
        return re.sub(r"\s+", " ", lowered)

    base = (native_text or "").strip()
    parts: List[str] = [base] if base else []
    seen_norm: List[str] = [_norm(base)] if base else []
    ocr = (ocr_text or "").strip()
    if ocr:
        ocr = ocr[:OCR_APPEND_MAX_CHARS].strip()
        ocr_norm = _norm(ocr)
        is_duplicate = any(
            existing and (ocr_norm == existing or ocr_norm in existing or existing in ocr_norm)
            for existing in seen_norm
        )
        if not is_duplicate:
            parts.append(f"[OCR]\n{ocr}")
            seen_norm.append(ocr_norm)
    caption = (vision_caption or "").strip()
    if caption:
        caption = caption[:VISION_APPEND_MAX_CHARS].strip()
        cap_norm = _norm(caption)
        is_duplicate = any(
            existing and (cap_norm == existing or cap_norm in existing or existing in cap_norm)
            for existing in seen_norm
        )
        if not is_duplicate:
            parts.append(f"[VISION]\n{caption}")
    merged = "\n\n".join(parts).strip()
    if len(merged) > MERGED_CHUNK_MAX_CHARS:
        merged = merged[:MERGED_CHUNK_MAX_CHARS].rstrip()
    return merged or base


def _content_sha256_hex(file_bytes: bytes) -> str:
    return hashlib.sha256(file_bytes).hexdigest()


async def _delete_chunk_points(collection_name: str, parent_file_id: str) -> int:
    """Delete all chunk points for a parent file before reindexing."""
    chunk_filter = Filter(
        must=[
            FieldCondition(key="parent_file_id", match=MatchValue(value=parent_file_id)),
            FieldCondition(key="record_type", match=MatchValue(value="chunk")),
        ]
    )
    try:
        points, _ = await vector_service.client.scroll(
            collection_name=collection_name,
            scroll_filter=chunk_filter,
            limit=5000,
            with_payload=False,
            with_vectors=False,
        )
    except Exception as exc:
        message = str(exc).lower()
        if "index required but not found" not in message:
            raise
        fallback_filter = Filter(
            must=[FieldCondition(key="parent_file_id", match=MatchValue(value=parent_file_id))]
        )
        points, _ = await vector_service.client.scroll(
            collection_name=collection_name,
            scroll_filter=fallback_filter,
            limit=5000,
            with_payload=True,
            with_vectors=False,
        )
        points = [
            point for point in points
            if (point.payload or {}).get("record_type") == "chunk"
        ]
    if not points:
        return 0
    await vector_service.client.delete(
        collection_name=collection_name,
        points_selector=[str(point.id) for point in points],
    )
    return len(points)


async def _mark_job_failed(
    *,
    collection_name: str,
    file_id: str,
    job_id: Optional[str],
    error_message: str,
) -> None:
    """Persist deterministic failed status for worker-level failures."""
    await vector_service.client.set_payload(
        collection_name=collection_name,
        payload={
            "ingestion_status": "failed",
            "ingestion_error": error_message,
            "ingestion_job_status": "failed" if job_id else None,
            "ingestion_job_error_message": error_message if job_id else None,
            "ingestion_job_completed_at": datetime.now().isoformat() if job_id else None,
        },
        points=[file_id],
    )


async def mark_ingestion_job_permanently_failed(
    *,
    collection_name: str,
    file_id: str,
    job_id: Optional[str],
    reason: str,
    error_message: str,
    attempt_count: Optional[int] = None,
    prior_failure_reason: Optional[str] = None,
    prior_error_message: Optional[str] = None,
) -> None:
    """Persist deterministic terminal status for permanent ingestion failures."""
    await vector_service.client.set_payload(
        collection_name=collection_name,
        payload={
            "ingestion_status": "permanently_failed",
            "ingestion_error": error_message,
            "ingestion_job_status": "permanently_failed" if job_id else None,
            "ingestion_job_error_message": error_message if job_id else None,
            "ingestion_job_completed_at": datetime.now().isoformat() if job_id else None,
            "ingestion_job_attempt_count": attempt_count if attempt_count is not None and job_id else None,
            "ingestion_failure_reason": reason,
            "ingestion_previous_failure_reason": prior_failure_reason,
            "ingestion_previous_error": prior_error_message,
        },
        points=[file_id],
    )


def _parse_payload_datetime(value: Any) -> Optional[datetime]:
    raw = str(value or "").strip()
    if not raw:
        return None
    try:
        normalized = raw.replace("Z", "+00:00")
        parsed = datetime.fromisoformat(normalized)
        if parsed.tzinfo is None:
            return parsed.replace(tzinfo=timezone.utc)
        return parsed.astimezone(timezone.utc)
    except Exception:
        return None


def _build_worker_payload(job_id: str, file_id: str, collection_name: str) -> Dict[str, str]:
    return {
        "job_id": job_id,
        "file_id": file_id,
        "collection_name": collection_name,
    }


async def _enqueue_cloud_task(job_id: str, file_id: str, collection_name: str) -> None:
    """Create a durable Cloud Task for ingestion worker execution."""
    settings = get_settings()
    if not settings.INGESTION_USE_CLOUD_TASKS:
        return
    if not settings.GCP_PROJECT_ID or not settings.INGESTION_WORKER_URL:
        raise RuntimeError("Cloud Tasks ingestion is enabled but GCP_PROJECT_ID/INGESTION_WORKER_URL is missing")

    payload = _build_worker_payload(job_id, file_id, collection_name)

    def _create_task_sync() -> None:
        client = tasks_v2.CloudTasksClient()
        parent = client.queue_path(
            settings.GCP_PROJECT_ID,
            settings.INGESTION_TASKS_LOCATION,
            settings.INGESTION_TASKS_QUEUE,
        )
        task_name = f"{parent}/tasks/{job_id}"
        headers = {"Content-Type": "application/json"}
        if settings.INGESTION_WORKER_TOKEN:
            headers["X-Ingestion-Worker-Token"] = settings.INGESTION_WORKER_TOKEN

        http_request: Dict[str, Any] = {
            "http_method": tasks_v2.HttpMethod.POST,
            "url": settings.INGESTION_WORKER_URL,
            "headers": headers,
            "body": json.dumps(payload).encode("utf-8"),
        }

        if settings.INGESTION_TASKS_SERVICE_ACCOUNT_EMAIL:
            http_request["oidc_token"] = {
                "service_account_email": settings.INGESTION_TASKS_SERVICE_ACCOUNT_EMAIL,
                "audience": settings.INGESTION_WORKER_URL,
            }

        task = {"name": task_name, "http_request": http_request}
        try:
            client.create_task(request={"parent": parent, "task": task})
        except AlreadyExists:
            # Task already exists for this job_id; treat as idempotent enqueue.
            return

    await run_in_threadpool(_create_task_sync)


async def create_ingestion_job(
    *,
    collection_name: str,
    file_id: str,
    tenant_id: str,
    initial_status: str = "queued",
) -> str:
    """Create/attach a durable ingestion job record on the parent file payload."""
    job_id = str(uuid.uuid4())
    now = datetime.now().isoformat()
    await vector_service.client.set_payload(
        collection_name=collection_name,
        payload={
            "ingestion_job_id": job_id,
            "ingestion_job_status": initial_status,
            "ingestion_job_attempt_count": 0,
            "ingestion_job_created_at": now,
            "ingestion_job_started_at": None,
            "ingestion_job_completed_at": None,
            "ingestion_job_error_message": None,
            "ingestion_job_processing_duration_ms": None,
            "ingestion_status": initial_status,
            "ingestion_queue_provider": "cloud_tasks" if get_settings().INGESTION_USE_CLOUD_TASKS else "inline",
            "tenant_scope": tenant_id,
        },
        points=[file_id],
    )
    return job_id


async def enqueue_ingestion_job(
    *,
    collection_name: str,
    file_id: str,
    tenant_id: str,
) -> str:
    """Create job metadata and enqueue durable worker task."""
    job_id = await create_ingestion_job(
        collection_name=collection_name,
        file_id=file_id,
        tenant_id=tenant_id,
        initial_status="queued",
    )
    await _enqueue_cloud_task(job_id, file_id, collection_name)
    return job_id


async def _run_ingestion_pipeline(
    *,
    job_id: Optional[str],
    collection_name: str,
    point_id: str,
    file_bytes: bytes,
    filename: str,
    file_type: str,
    tenant_id: str,
    is_global: bool,
    tags: List[str],
    file_url: str,
    size_str: str,
    file_size: int,
    upload_date: str,
    preview_url: Optional[str],
    preview_type: Optional[str],
    preview_status: str,
    preview_error: Optional[str],
    initial_ai_enabled: bool,
    source_content_sha256: str,
    attempt_count: int,
) -> None:
    """Async extraction -> quality gate -> chunk embedding pipeline."""
    settings = get_settings()
    ai_enabled_for_retrieval = bool(initial_ai_enabled)
    started = time.perf_counter()

    async def _abort_if_ai_disabled(stage: str) -> bool:
        latest_points = await vector_service.client.retrieve(
            collection_name=collection_name,
            ids=[point_id],
            with_payload=True,
            with_vectors=False,
        )
        latest_payload = latest_points[0].payload if latest_points else {}
        if bool((latest_payload or {}).get("ai_enabled", False)):
            return False
        now_iso = datetime.now().isoformat()
        current_status = str((latest_payload or {}).get("ingestion_status") or "uploaded").strip().lower()
        await vector_service.client.set_payload(
            collection_name=collection_name,
            payload={
                "ingestion_status": "uploaded" if current_status in {"queued", "processing", "uploaded"} else current_status,
                "ingestion_job_status": "cancelled" if job_id else None,
                "ingestion_job_completed_at": now_iso if job_id else None,
                "ingestion_job_error_message": f"Skipped during {stage}: Riley Memory disabled." if job_id else None,
            },
            points=[point_id],
        )
        return True

    try:
        started_at = datetime.now().isoformat()
        if await _abort_if_ai_disabled("pre_extract"):
            return
        await vector_service.client.set_payload(
            collection_name=collection_name,
            payload={
                "ingestion_status": "processing",
                "ingestion_started_at": started_at,
                "ingestion_job_status": "processing" if job_id else None,
                "ingestion_job_started_at": started_at if job_id else None,
                "ingestion_job_attempt_count": attempt_count if job_id else None,
                "ingestion_job_error_message": None if job_id else None,
            },
            points=[point_id],
        )

        raw_text = await extract_text(file_bytes, filename)
        structured_segments = await _extract_structured_segments(file_bytes, filename, raw_text)
        merged_cleaned_text, prepared_segments = _prepare_segments_for_chunking(structured_segments)
        cleaned_text = _clean_extracted_text(merged_cleaned_text or raw_text)[:MAX_CHARS_TOTAL]
        raw_text_capped = (raw_text or "")[:MAX_CHARS_TOTAL]
        quality_status, quality_reason = _assess_extraction_quality(cleaned_text, filename, file_size)

        ocr_enabled = False
        ocr_status = "not_requested"
        ocr_text_present = False
        ocr_confidence: Optional[float] = None
        multimodal_status = "native_only"

        if await _abort_if_ai_disabled("pre_ocr"):
            return

        if settings.OCR_ENABLED_SYSTEMWIDE and _ocr_required_for_file(
            file_type=file_type,
            quality_status=quality_status,
            cleaned_text=cleaned_text,
            file_size=file_size,
        ):
            ocr_enabled = True
            ocr_status = "processing"
            multimodal_status = "ocr_attempted"
            try:
                ocr_result: Dict[str, Any] = {"text": "", "confidence": None, "pages": []}
                file_ext = (file_type or "").lower()
                if file_ext in IMAGE_EXTENSIONS:
                    ocr_result = await run_ocr(file_bytes, max_chars=settings.OCR_MAX_CHARS)
                elif file_ext == "pdf":
                    source_uri = gcs_uri_from_url(file_url, settings.GCS_BUCKET_NAME)
                    if source_uri:
                        ocr_result = await run_ocr_document_from_gcs(
                            gcs_source_uri=source_uri,
                            output_bucket=settings.GCS_BUCKET_NAME,
                            output_prefix="ocr-results",
                            mime_type="application/pdf",
                            max_chars=max(MAX_CHARS_TOTAL, settings.OCR_MAX_CHARS * 2),
                        )
                elif file_ext in {"pptx", "ppt"}:
                    # PPT/PPTX previews are stored as PDF at deterministic path.
                    preview_source_uri = f"gs://{settings.GCS_BUCKET_NAME}/{settings.PREVIEW_BUCKET_PATH_PREFIX}/{point_id}.pdf"
                    if preview_status == "complete":
                        ocr_result = await run_ocr_document_from_gcs(
                            gcs_source_uri=preview_source_uri,
                            output_bucket=settings.GCS_BUCKET_NAME,
                            output_prefix="ocr-results",
                            mime_type="application/pdf",
                            max_chars=max(MAX_CHARS_TOTAL, settings.OCR_MAX_CHARS * 2),
                        )

                ocr_segments = _build_ocr_segments(
                    filename=filename,
                    file_type=file_type,
                    ocr_result=ocr_result,
                )
                if ocr_segments:
                    ocr_text_present = True
                    conf_raw = ocr_result.get("confidence")
                    if isinstance(conf_raw, (float, int)):
                        ocr_confidence = max(0.0, min(1.0, float(conf_raw)))
                    ocr_status = "complete"
                    multimodal_status = "ocr_enriched"
                    # If native text is weak, prioritize OCR segments; otherwise append.
                    if quality_status != "indexed":
                        structured_segments = ocr_segments
                    else:
                        structured_segments = [*structured_segments, *ocr_segments]
                    merged_cleaned_text, prepared_segments = _prepare_segments_for_chunking(structured_segments)
                    cleaned_text = _clean_extracted_text(merged_cleaned_text or raw_text)[:MAX_CHARS_TOTAL]
                    quality_status, quality_reason = _assess_extraction_quality(cleaned_text, filename, file_size)
                    if file_ext in IMAGE_EXTENSIONS:
                        quality_status = "indexed"
                        quality_reason = "OCR extracted text from image asset"
                    if quality_status == "ocr_needed":
                        quality_status = "indexed"
                        quality_reason = "OCR recovered text from visual content"
                else:
                    ocr_status = "failed"
                    multimodal_status = "ocr_unavailable"
            except Exception as ocr_exc:
                ocr_status = "failed"
                multimodal_status = "ocr_failed"
                logger.warning(
                    "ocr_multimodal_failed file_id=%s file_type=%s error=%s",
                    point_id,
                    file_type,
                    ocr_exc,
                )

        vision_enabled = False
        vision_status = "not_requested"
        vision_segments_annotated = 0
        if settings.RILEY_VISION_ENABLED and settings.OPENAI_API_KEY:
            max_vision_segments = max(1, int(settings.RILEY_VISION_MAX_SEGMENTS))
            candidate_pages = _vision_candidate_pages(
                file_type=file_type,
                quality_status=quality_status,
                segments=structured_segments,
                ocr_enabled=ocr_enabled,
                max_pages=max_vision_segments,
            )
            if candidate_pages:
                vision_enabled = True
                vision_status = "processing"
                rendered_pages: Dict[int, bytes] = {}
                file_ext = (file_type or "").lower()
                if file_ext in IMAGE_EXTENSIONS:
                    rendered_pages = {1: file_bytes}
                elif file_ext == "pdf":
                    rendered_pages = await _render_pdf_pages_for_vision(
                        pdf_bytes=file_bytes,
                        page_numbers=candidate_pages,
                        max_pages=max_vision_segments,
                    )
                elif file_ext in {"pptx", "ppt"} and preview_url:
                    try:
                        preview_bytes = await StorageService.download_file(preview_url)
                        rendered_pages = await _render_pdf_pages_for_vision(
                            pdf_bytes=preview_bytes,
                            page_numbers=candidate_pages,
                            max_pages=max_vision_segments,
                        )
                    except Exception as preview_exc:
                        logger.warning(
                            "vision_preview_render_failed file_id=%s error=%s",
                            point_id,
                            type(preview_exc).__name__,
                        )
                repeated_failure_counts: Dict[str, int] = {}
                max_repeated_page_failures = 3
                for segment in structured_segments:
                    if vision_segments_annotated >= max_vision_segments:
                        break
                    loc_val = str(segment.get("location_value") or "").strip()
                    page_num = int(loc_val) if loc_val.isdigit() else None
                    if page_num is None:
                        continue
                    page_bytes = rendered_pages.get(page_num)
                    if not page_bytes:
                        continue
                    try:
                        vision_result = await asyncio.wait_for(
                            summarize_visual_content(
                                image_bytes=page_bytes,
                                model=settings.RILEY_VISION_MODEL,
                                timeout_seconds=settings.RILEY_VISION_TIMEOUT_SECONDS,
                            ),
                            timeout=max(5, int(settings.RILEY_VISION_TIMEOUT_SECONDS) + 2),
                        )
                        segment["vision_caption"] = str(vision_result.get("vision_caption") or "")
                        segment["has_visual_content"] = bool(vision_result.get("has_visual_content"))
                        segment["visual_type"] = vision_result.get("visual_type")
                        segment["multimodal_status"] = "vision_enriched"
                        vision_segments_annotated += 1
                    except Exception as vision_exc:
                        error_type = type(vision_exc).__name__
                        error_message = str(vision_exc).strip() or "<empty>"
                        failure_signature = f"{error_type}:{error_message}"
                        repeated_failure_counts[failure_signature] = (
                            repeated_failure_counts.get(failure_signature, 0) + 1
                        )
                        repeat_count = repeated_failure_counts[failure_signature]
                        logger.warning(
                            "vision_caption_failed file_id=%s filename=%s page=%s error_type=%s error_message=%s repeat_count=%s",
                            point_id,
                            filename,
                            page_num,
                            error_type,
                            error_message,
                            repeat_count,
                        )
                        if repeat_count >= max_repeated_page_failures:
                            logger.warning(
                                "vision_caption_short_circuit file_id=%s filename=%s page=%s failure_signature=%s threshold=%s",
                                point_id,
                                filename,
                                page_num,
                                failure_signature,
                                max_repeated_page_failures,
                            )
                            break
                if vision_segments_annotated > 0:
                    vision_status = "complete"
                    multimodal_status = "vision_enriched"
                    # Rebuild merged text so vision-only or weak-native docs remain searchable.
                    for segment in structured_segments:
                        caption = str(segment.get("vision_caption") or "").strip()
                        if not caption:
                            continue
                        base_text = str(segment.get("text") or segment.get("raw_text") or "").strip()
                        if caption and caption not in base_text:
                            enriched = f"{base_text}\n\n[VISION]\n{caption}".strip() if base_text else f"[VISION]\n{caption}"
                            segment["text"] = enriched
                            segment["raw_text"] = enriched
                    merged_cleaned_text, prepared_segments = _prepare_segments_for_chunking(structured_segments)
                    cleaned_text = _clean_extracted_text(merged_cleaned_text or raw_text)[:MAX_CHARS_TOTAL]
                else:
                    vision_status = "failed"

        content_signals = _compute_content_signals(
            native_text=raw_text,
            merged_text=cleaned_text,
            segments=structured_segments,
        )
        final_status, decision_code, decision_why = _decide_final_ingestion_status(
            file_type=file_type,
            quality_status=quality_status,
            quality_reason=quality_reason,
            signals=content_signals,
        )
        quality_status = final_status
        quality_reason = decision_why
        pptx_slide_stats: Dict[str, int] = {"total_slides": 0, "indexed_slides": 0, "skipped_slides": 0}
        if file_type == "pptx":
            _, pptx_slide_stats = _build_pptx_slide_micro_chunks(
                prepared_segments,
                chunk_size_tokens=MICRO_CHUNK_SIZE_TOKENS,
                overlap_tokens=MICRO_CHUNK_OVERLAP_TOKENS,
                max_chars_per_chunk=MERGED_CHUNK_MAX_CHARS,
            )

        if quality_status == "low_text":
            deleted_points = await _delete_chunk_points(collection_name, point_id)
            if deleted_points > 0:
                logger.info(
                    "vectors_deleted_before_reindex file_id=%s collection=%s deleted_vectors=%s",
                    point_id,
                    collection_name,
                    deleted_points,
                )
            duration_ms = int((time.perf_counter() - started) * 1000)
            completed_at = datetime.now().isoformat()
            logger.info(
                "ingestion_decision file_id=%s filename=%s file_type=%s native_char_count=%s ocr_char_count=%s "
                "vision_caption_count=%s merged_char_count=%s nonempty_page_count=%s total_page_count=%s final_status=%s decision_reason=%s",
                point_id,
                filename,
                file_type,
                content_signals.get("native_char_count", 0),
                content_signals.get("ocr_char_count", 0),
                content_signals.get("vision_caption_count", 0),
                content_signals.get("merged_char_count", 0),
                content_signals.get("nonempty_page_count", 0),
                content_signals.get("total_page_count", 0),
                quality_status,
                quality_reason,
            )
            low_text_payload: Dict[str, Any] = {
                    "record_type": "file",
                    "filename": filename,
                    "file_type": file_type,
                    "type": file_type,
                    "url": file_url,
                    "is_global": is_global,
                    "client_id": None if is_global else tenant_id,
                    "tags": tags,
                    "size": size_str,
                    "size_bytes": file_size,
                    "upload_date": upload_date,
                    "uploaded_at": upload_date,
                    "preview_url": preview_url,
                    "preview_type": preview_type,
                    "preview_status": preview_status,
                    "preview_error": preview_error,
                    "raw_content": raw_text_capped,
                    "cleaned_content": cleaned_text[:MAX_CHARS_TOTAL],
                    "content": cleaned_text[:MAX_CHARS_TOTAL],
                    "content_preview": cleaned_text[:CONTENT_PREVIEW_LENGTH],
                    "ai_enabled": False,
                    "ingestion_status": quality_status,
                    "ingestion_error": quality_reason,
                    "extracted_char_count": len(cleaned_text),
                    "chunk_count": 0,
                    "embedding_model": settings.EMBEDDING_MODEL,
                    "embedding_tokens_estimate": 0,
                    "embedding_cost_estimate_usd": 0.0,
                    "bm25_enabled": False,
                    "chunk_profiles": {"micro": 0, "macro": 0},
                    "processing_duration_ms": duration_ms,
                    "ocr_enabled": ocr_enabled,
                    "ocr_status": ocr_status if ocr_enabled else ("not_requested" if file_type in IMAGE_EXTENSIONS else None),
                    "ocr_text_present": ocr_text_present,
                    "ocr_confidence": ocr_confidence,
                    "multimodal_status": multimodal_status,
                    "multimodal_enabled": bool(ocr_enabled or vision_enabled),
                    "ocr_processed": bool(ocr_enabled),
                    "vision_processed": bool(vision_enabled),
                    "visual_chunk_count": 0,
                    "indexed_content_sha256": None,
                    "vision_enabled": vision_enabled,
                    "vision_status": vision_status if vision_enabled else "not_requested",
                    "vision_segments_annotated": vision_segments_annotated,
                    "ingestion_job_status": quality_status if job_id else None,
                    "ingestion_job_completed_at": completed_at if job_id else None,
                    "ingestion_job_error_message": quality_reason if job_id else None,
                    "ingestion_job_processing_duration_ms": duration_ms if job_id else None,
                    "analysis_status": "not_requested",
                    "analysis_completed_at": None,
                    "analysis_error": None,
                    "ingestion_decision_code": decision_code,
                    "ingestion_decision_why": quality_reason,
                    "ingestion_decision": {
                        "code": decision_code,
                        "why": quality_reason,
                        "native_char_count": content_signals.get("native_char_count", 0),
                        "ocr_char_count": content_signals.get("ocr_char_count", 0),
                        "vision_caption_count": content_signals.get("vision_caption_count", 0),
                        "merged_char_count": content_signals.get("merged_char_count", 0),
                        "nonempty_page_count": content_signals.get("nonempty_page_count", 0),
                        "total_page_count": content_signals.get("total_page_count", 0),
                    },
                }
            if file_type == "pptx":
                low_text_payload.update(
                    {
                        "pptx_total_slides": int(pptx_slide_stats.get("total_slides") or 0),
                        "pptx_indexed_slide_count": int(pptx_slide_stats.get("indexed_slides") or 0),
                        "pptx_skipped_slide_count": int(pptx_slide_stats.get("skipped_slides") or 0),
                    }
                )
            await vector_service.client.set_payload(
                collection_name=collection_name,
                payload=low_text_payload,
                points=[point_id],
            )
            logger.info(
                "ingestion_complete file_id=%s tenant=%s status=%s chars=%s chunks=0 reason=%s duration_ms=%s",
                point_id, tenant_id, quality_status, len(cleaned_text), quality_reason, duration_ms
            )
            return

        latest_parent_points = await vector_service.client.retrieve(
            collection_name=collection_name,
            ids=[point_id],
            with_payload=True,
            with_vectors=False,
        )
        if latest_parent_points:
            latest_parent_payload = latest_parent_points[0].payload or {}
            ai_enabled_for_retrieval = bool(latest_parent_payload.get("ai_enabled", ai_enabled_for_retrieval))
        if await _abort_if_ai_disabled("pre_chunking"):
            return

        pptx_chunk_stats: Dict[str, int] = {"total_slides": 0, "indexed_slides": 0, "skipped_slides": 0}
        micro_chunks: List[Dict[str, Any]] = []
        if file_type == "pptx":
            micro_chunks, pptx_chunk_stats = _build_pptx_slide_micro_chunks(
                prepared_segments,
                chunk_size_tokens=MICRO_CHUNK_SIZE_TOKENS,
                overlap_tokens=MICRO_CHUNK_OVERLAP_TOKENS,
                max_chars_per_chunk=MERGED_CHUNK_MAX_CHARS,
            )
        if not micro_chunks:
            if file_type == "pptx":
                fallback_text = cleaned_text[:MAX_CHARS_TOTAL].strip()
                fallback_chunks = _chunk_text_with_offsets(
                    fallback_text,
                    chunk_size_tokens=MICRO_CHUNK_SIZE_TOKENS,
                    overlap_tokens=MICRO_CHUNK_OVERLAP_TOKENS,
                )
                if not fallback_chunks and fallback_text:
                    fallback_chunks = [{"text": fallback_text, "char_start": 0, "char_end": len(fallback_text)}]
                for idx, fallback_chunk in enumerate(fallback_chunks):
                    micro_chunks.append(
                        {
                            "text": str(fallback_chunk.get("text") or ""),
                            "char_start": int(fallback_chunk.get("char_start") or 0),
                            "char_end": int(fallback_chunk.get("char_end") or 0),
                            "slide_number": 1,
                            "slide_title": None,
                            "slide_subchunk_index": idx,
                        }
                    )
                if micro_chunks:
                    inferred_total = max(1, int(pptx_chunk_stats.get("total_slides") or 0))
                    pptx_chunk_stats = {
                        "total_slides": inferred_total,
                        "indexed_slides": 1,
                        "skipped_slides": max(0, inferred_total - 1),
                    }
            else:
                micro_chunks = _chunk_text_with_offsets(
                    cleaned_text,
                    chunk_size_tokens=MICRO_CHUNK_SIZE_TOKENS,
                    overlap_tokens=MICRO_CHUNK_OVERLAP_TOKENS,
                )
        macro_chunks: List[Dict[str, Any]] = []
        if file_type != "pptx":
            macro_chunks = _chunk_text_with_offsets(
                cleaned_text,
                chunk_size_tokens=MACRO_CHUNK_SIZE_TOKENS,
                overlap_tokens=MACRO_CHUNK_OVERLAP_TOKENS,
            )
        if not micro_chunks:
            fallback_text = cleaned_text[:MAX_CHARS_TOTAL]
            if file_type == "pptx":
                if fallback_text.strip():
                    micro_chunks = [
                        {
                            "text": fallback_text,
                            "char_start": 0,
                            "char_end": len(fallback_text),
                            "slide_number": 1,
                            "slide_title": None,
                            "slide_subchunk_index": 0,
                        }
                    ]
            else:
                micro_chunks = [{"text": fallback_text, "char_start": 0, "char_end": len(fallback_text)}]
        if file_type != "pptx" and not macro_chunks:
            macro_chunks = [{"text": cleaned_text[:MAX_CHARS_TOTAL], "char_start": 0, "char_end": len(cleaned_text[:MAX_CHARS_TOTAL])}]

        deleted_points = await _delete_chunk_points(collection_name, point_id)
        if deleted_points > 0:
            logger.info(
                "vectors_deleted_before_reindex file_id=%s collection=%s deleted_vectors=%s",
                point_id,
                collection_name,
                deleted_points,
            )

        chunk_rows: List[Dict[str, Any]] = []
        total_tokens = 0
        bm25_available = await vector_service.bm25_enabled_for_collection(collection_name)
        chunk_processed_count = 0
        for chunk_type, chunk_list in (("micro", micro_chunks), ("macro", macro_chunks)):
            for chunk_idx, chunk in enumerate(chunk_list):
                if chunk_processed_count % 5 == 0:
                    if await _abort_if_ai_disabled("chunk_embedding"):
                        return
                chunk_text = str(chunk.get("text") or "").strip()
                if not chunk_text:
                    continue
                char_start = int(chunk.get("char_start") or 0)
                char_end = int(chunk.get("char_end") or len(chunk_text))
                best_segment = _best_segment_for_span(char_start, char_end, prepared_segments)
                native_chunk_text = chunk_text
                ocr_chunk_text = (
                    str(best_segment.get("raw_text") or "").strip()
                    if bool(best_segment.get("ocr_text_present"))
                    else None
                )
                vision_caption = str(best_segment.get("vision_caption") or "").strip() or None
                merged_chunk_text = _merge_chunk_text(
                    native_text=native_chunk_text,
                    ocr_text=ocr_chunk_text,
                    vision_caption=vision_caption,
                )
                vector = await _generate_embedding(merged_chunk_text)
                token_estimate = _estimate_tokens(merged_chunk_text)
                total_tokens += token_estimate
                chunk_processed_count += 1
                point_uuid = _chunk_point_uuid(point_id, chunk_type, chunk_idx)
                chunk_id = f"{point_id}::chunk::{chunk_type}::{chunk_idx}"
                chunk_payload: Dict[str, Any] = {
                    "record_type": "chunk",
                    "parent_file_id": point_id,
                    "chunk_id": chunk_id,
                    "chunk_type": chunk_type,
                    "chunk_index": chunk_idx,
                    "filename": filename,
                    "file_type": file_type,
                    "type": file_type,
                    "client_id": None if is_global else tenant_id,
                    "is_global": is_global,
                    "tenant_id": tenant_id,
                    "url": file_url,
                    "tags": tags,
                    "ai_enabled": ai_enabled_for_retrieval,
                    "ingestion_status": quality_status,
                    "text": merged_chunk_text,
                    "content": merged_chunk_text,
                    "raw_text": str(best_segment.get("raw_text") or chunk_text),
                    "content_preview": merged_chunk_text[:CONTENT_PREVIEW_LENGTH],
                    "upload_date": upload_date,
                    "uploaded_at": upload_date,
                    "chunk_token_estimate": token_estimate,
                    "location_type": str(best_segment.get("location_type") or "unknown"),
                    "location_value": str(best_segment.get("location_value") or ""),
                    "section_path": str(best_segment.get("section_path") or ""),
                    "char_start": char_start,
                    "char_end": char_end,
                    "ocr_enabled": ocr_enabled,
                    "ocr_status": str(best_segment.get("ocr_status") or ocr_status),
                    "ocr_text_present": bool(best_segment.get("ocr_text_present") or ocr_text_present),
                    "ocr_confidence": (
                        float(best_segment.get("ocr_confidence"))
                        if isinstance(best_segment.get("ocr_confidence"), (float, int))
                        else ocr_confidence
                    ),
                    "ocr_text": (
                        str(best_segment.get("raw_text") or "")[:MAX_CHARS_TOTAL]
                        if bool(best_segment.get("ocr_text_present"))
                        else None
                    ),
                    "multimodal_status": str(best_segment.get("multimodal_status") or multimodal_status),
                    "native_text": native_chunk_text,
                    "vision_caption": vision_caption,
                    "has_visual_content": bool(best_segment.get("has_visual_content")),
                    "visual_type": best_segment.get("visual_type"),
                    "vision_enabled": vision_enabled,
                    "vision_status": vision_status if vision_enabled else "not_requested",
                }
                location_type = str(chunk_payload.get("location_type") or "").strip().lower()
                location_value = str(chunk_payload.get("location_value") or "").strip()
                slide_number: Optional[int] = None
                if location_type == "slide" and location_value.isdigit():
                    slide_number = int(location_value)
                slide_title = str(best_segment.get("slide_title") or chunk.get("slide_title") or "").strip() or None
                if file_type == "pptx":
                    slide_number = int(chunk.get("slide_number")) if isinstance(chunk.get("slide_number"), int) else slide_number
                    slide_subchunk_index = (
                        int(chunk.get("slide_subchunk_index"))
                        if isinstance(chunk.get("slide_subchunk_index"), int)
                        else 0
                    )
                    chunk_payload["content_type"] = "presentation_slide"
                    chunk_payload["slide_number"] = slide_number
                    chunk_payload["slide_title"] = slide_title
                    chunk_payload["slide_subchunk_index"] = slide_subchunk_index
                    if slide_number is not None:
                        chunk_payload["location_type"] = "slide"
                        chunk_payload["location_value"] = str(slide_number)
                        chunk_payload["section_path"] = f"slide:{slide_number}"
                chunk_rows.append(
                    {
                        "id": point_uuid,
                        "dense_vector": vector,
                        "chunk_text": merged_chunk_text,
                        "payload": chunk_payload,
                    }
                )

        def _build_chunk_points(include_bm25: bool) -> List[PointStruct]:
            points: List[PointStruct] = []
            for row in chunk_rows:
                payload = dict(row["payload"])
                payload["bm25_enabled"] = include_bm25
                vector_payload: Any = row["dense_vector"]
                if include_bm25:
                    vector_payload = {
                        "": row["dense_vector"],
                        "bm25": qdrant_models.Document(text=row["chunk_text"], model=BM25_MODEL),
                    }
                points.append(
                    PointStruct(
                        id=row["id"],
                        vector=vector_payload,
                        payload=payload,
                    )
                )
            return points

        if await _abort_if_ai_disabled("pre_chunk_persist"):
            return

        if bm25_available:
            try:
                await vector_service.client.upsert(
                    collection_name=collection_name,
                    points=_build_chunk_points(include_bm25=True),
                )
            except Exception as bm25_exc:
                if vector_service.is_missing_bm25_vector_error(bm25_exc):
                    # Try one safe migration refresh in case the collection predates sparse config.
                    bm25_available = await vector_service.refresh_bm25_support(collection_name)
                    if bm25_available:
                        try:
                            await vector_service.client.upsert(
                                collection_name=collection_name,
                                points=_build_chunk_points(include_bm25=True),
                            )
                        except Exception as bm25_retry_exc:
                            bm25_available = False
                            vector_service.mark_bm25_unavailable(collection_name, str(bm25_retry_exc))
                    else:
                        vector_service.mark_bm25_unavailable(collection_name, str(bm25_exc))
                else:
                    logger.warning(
                        "bm25_chunk_upsert_failed collection=%s file_id=%s error=%s; falling back to dense-only indexing",
                        collection_name,
                        point_id,
                        bm25_exc,
                    )
                    bm25_available = False
            if not bm25_available:
                await vector_service.client.upsert(
                    collection_name=collection_name,
                    points=_build_chunk_points(include_bm25=False),
                )
        else:
            await vector_service.client.upsert(
                collection_name=collection_name,
                points=_build_chunk_points(include_bm25=False),
            )

        visual_chunk_count = sum(
            1 for row in chunk_rows if bool((row.get("payload") or {}).get("has_visual_content"))
        )
        duration_ms = int((time.perf_counter() - started) * 1000)
        embedding_cost_estimate = round((total_tokens / 1000.0) * EMBEDDING_COST_PER_1K_TOKENS_USD, 6)
        logger.info(
            "ingestion_decision file_id=%s filename=%s file_type=%s native_char_count=%s ocr_char_count=%s "
            "vision_caption_count=%s merged_char_count=%s nonempty_page_count=%s total_page_count=%s final_status=%s decision_reason=%s",
            point_id,
            filename,
            file_type,
            content_signals.get("native_char_count", 0),
            content_signals.get("ocr_char_count", 0),
            content_signals.get("vision_caption_count", 0),
            content_signals.get("merged_char_count", 0),
            content_signals.get("nonempty_page_count", 0),
            content_signals.get("total_page_count", 0),
            quality_status,
            quality_reason,
        )
        parent_payload: Dict[str, Any] = {
            "record_type": "file",
            "filename": filename,
            "file_type": file_type,
            "type": file_type,
            "url": file_url,
            "is_global": is_global,
            "tags": tags,
            "raw_content": raw_text_capped,
            "cleaned_content": cleaned_text[:MAX_CHARS_TOTAL],
            "content": cleaned_text[:MAX_CHARS_TOTAL],
            "content_preview": cleaned_text[:CONTENT_PREVIEW_LENGTH],
            "size": size_str,
            "size_bytes": file_size,
            "upload_date": upload_date,
            "uploaded_at": upload_date,
            "preview_url": preview_url,
            "preview_type": preview_type,
            "preview_status": preview_status,
            "preview_error": preview_error,
            "ai_enabled": ai_enabled_for_retrieval,
            "ingestion_status": quality_status,
            "ingestion_error": None,
            "extracted_char_count": len(cleaned_text),
            "chunk_count": len(chunk_rows),
            "embedding_model": settings.EMBEDDING_MODEL,
            "embedding_tokens_estimate": total_tokens,
            "embedding_cost_estimate_usd": embedding_cost_estimate,
            "bm25_enabled": bm25_available,
            "chunk_profiles": {
                "micro": len(micro_chunks),
                "macro": len(macro_chunks),
            },
            "indexed_content_sha256": source_content_sha256,
            "processing_duration_ms": duration_ms,
            "ocr_enabled": ocr_enabled,
            "ocr_status": ocr_status if ocr_enabled else ("not_requested" if file_type in IMAGE_EXTENSIONS else None),
            "ocr_text_present": ocr_text_present,
            "ocr_confidence": ocr_confidence,
            "multimodal_status": multimodal_status,
            "multimodal_enabled": bool(ocr_enabled or vision_enabled),
            "ocr_processed": bool(ocr_enabled),
            "vision_processed": bool(vision_enabled),
            "visual_chunk_count": visual_chunk_count,
            "vision_enabled": vision_enabled,
            "vision_status": vision_status if vision_enabled else "not_requested",
            "vision_segments_annotated": vision_segments_annotated,
            "analysis_status": "queued" if settings.RILEY_DOC_INTEL_ENABLED else "not_requested",
            "analysis_completed_at": None,
            "analysis_error": None,
            "ingestion_decision_code": decision_code,
            "ingestion_decision_why": quality_reason,
            "ingestion_decision": {
                "code": decision_code,
                "why": quality_reason,
                "native_char_count": content_signals.get("native_char_count", 0),
                "ocr_char_count": content_signals.get("ocr_char_count", 0),
                "vision_caption_count": content_signals.get("vision_caption_count", 0),
                "merged_char_count": content_signals.get("merged_char_count", 0),
                "nonempty_page_count": content_signals.get("nonempty_page_count", 0),
                "total_page_count": content_signals.get("total_page_count", 0),
            },
        }
        if file_type == "pptx":
            parent_payload.update(
                {
                    "pptx_total_slides": int(pptx_chunk_stats.get("total_slides") or 0),
                    "pptx_indexed_slide_count": int(pptx_chunk_stats.get("indexed_slides") or 0),
                    "pptx_skipped_slide_count": int(pptx_chunk_stats.get("skipped_slides") or 0),
                }
            )
        if not is_global:
            parent_payload["client_id"] = tenant_id
        completed_at = datetime.now().isoformat()
        if await _abort_if_ai_disabled("pre_parent_persist"):
            return
        await vector_service.client.set_payload(
            collection_name=collection_name,
            payload=parent_payload,
            points=[point_id],
        )
        if job_id:
            await vector_service.client.set_payload(
                collection_name=collection_name,
                payload={
                    "ingestion_job_status": quality_status,
                    "ingestion_job_completed_at": completed_at,
                    "ingestion_job_error_message": None,
                    "ingestion_job_processing_duration_ms": duration_ms,
                },
                points=[point_id],
            )
        if settings.RILEY_DOC_INTEL_ENABLED:
            try:
                from app.services.document_intelligence import enqueue_document_intelligence_job

                await enqueue_document_intelligence_job(
                    collection_name=collection_name,
                    file_id=point_id,
                    tenant_id=tenant_id,
                    is_global=is_global,
                )
            except Exception as doc_intel_exc:
                logger.warning(
                    "doc_intel_enqueue_failed file_id=%s tenant=%s error=%s",
                    point_id,
                    tenant_id,
                    doc_intel_exc,
                )
        logger.info(
            "ingestion_complete file_id=%s tenant=%s status=%s chars=%s chunks=%s duration_ms=%s tokens=%s est_cost_usd=%s",
            point_id,
            tenant_id,
            quality_status,
            len(cleaned_text),
            len(chunk_rows),
            duration_ms,
            total_tokens,
            embedding_cost_estimate,
        )
    except Exception as exc:
        duration_ms = int((time.perf_counter() - started) * 1000)
        logger.exception("ingestion_failed file_id=%s tenant=%s error=%s", point_id, tenant_id, exc)
        fallback_signals = locals().get("content_signals") or {}
        logger.info(
            "ingestion_decision file_id=%s filename=%s file_type=%s native_char_count=%s ocr_char_count=%s "
            "vision_caption_count=%s merged_char_count=%s nonempty_page_count=%s total_page_count=%s final_status=%s decision_reason=%s",
            point_id,
            filename,
            file_type,
            fallback_signals.get("native_char_count", 0),
            fallback_signals.get("ocr_char_count", 0),
            fallback_signals.get("vision_caption_count", 0),
            fallback_signals.get("merged_char_count", 0),
            fallback_signals.get("nonempty_page_count", 0),
            fallback_signals.get("total_page_count", 0),
            "failed",
            f"failed: processing exception: {type(exc).__name__}",
        )
        completed_at = datetime.now().isoformat()
        await vector_service.client.set_payload(
            collection_name=collection_name,
            payload={
                "ingestion_status": "failed",
                "ingestion_error": str(exc),
                "processing_duration_ms": duration_ms,
                "ai_enabled": False,
                "bm25_enabled": False,
                "ingestion_job_status": "failed" if job_id else None,
                "ingestion_job_completed_at": completed_at if job_id else None,
                "ingestion_job_error_message": str(exc) if job_id else None,
                "ingestion_job_processing_duration_ms": duration_ms if job_id else None,
                "ocr_status": "failed",
                "multimodal_status": "ocr_failed",
                "multimodal_enabled": True,
                "ocr_processed": True,
                "vision_processed": False,
                "visual_chunk_count": 0,
                "analysis_status": "failed",
                "analysis_completed_at": completed_at,
                "analysis_error": "Ingestion failed; document intelligence was not generated.",
                "ingestion_decision_code": "failed_processing_exception",
                "ingestion_decision_why": f"failed: processing exception: {type(exc).__name__}",
                "ingestion_decision": {
                    "code": "failed_processing_exception",
                    "why": f"failed: processing exception: {type(exc).__name__}",
                },
            },
            points=[point_id],
        )


async def process_upload(
    file: UploadFile,
    tenant_id: str,
    tags: List[str],
    overwrite: bool = False,
) -> dict:
    """
    Process an uploaded file: upload to GCS, extract text, vectorize, and index in Qdrant.
    
    STANDARDIZED CONTENT STORAGE (SINGLE SOURCE OF TRUTH):
    - payload["content"]: FULL extracted text (capped at MAX_CHARS_TOTAL=20,000) for RAG
    - payload["content_preview"]: First 1000 chars for UI display
    - For images: content = "[Image file — OCR not requested]" (placeholder)
    
    This ensures RAG has access to full text, not just 1000 chars.
    
    Args:
        file: The uploaded file
        tenant_id: Tenant/client identifier for data isolation
        tags: List of tags associated with the file
        overwrite: If True, replace existing file with same name. If False, raise 409 on conflict.
        Note: OCR is NOT performed during upload. OCR must be explicitly
        requested via the dedicated POST /files/{file_id}/ocr endpoint.
        
    Returns:
        Dictionary containing:
            - id: The UUID of the indexed document
            - url: The public GCS URL to access the file
            - filename: The original filename
            - type: The file type/extension
    """
    settings = get_settings()

    if file.filename:
        file_extension = "." + file.filename.split(".")[-1] if "." in file.filename else ""
        base_name = file.filename.rsplit(".", 1)[0] if "." in file.filename else file.filename
        unique_filename = f"{uuid.uuid4()}_{base_name}{file_extension}"
    else:
        unique_filename = str(uuid.uuid4())

    try:
        file_content = await file.read()
        file_size = len(file_content)
        size_str = _format_file_size(file_size)
        upload_date = datetime.now().isoformat()
        filename = file.filename or "unknown"
        file_type = filename.split(".")[-1].lower() if "." in filename else "unknown"
        is_global_upload = tenant_id == "global"
        is_image = is_image_ext(filename)

        point_id = str(uuid.uuid4())

        # Upload original bytes first.
        file_url = await StorageService.upload_bytes(
            unique_filename,
            file_content,
            content_type=file.content_type or "application/octet-stream",
        )

        preview_url: Optional[str] = None
        preview_object_name: Optional[str] = None
        preview_type: Optional[str] = None
        preview_status: str = "not_requested"
        preview_error: Optional[str] = None
        if settings.ENABLE_PREVIEW_GENERATION and is_office_or_html(filename):
            preview_status = "processing"
            try:
                pdf_bytes = await generate_pdf_preview(file_content, filename)
                preview_object_name = f"{settings.PREVIEW_BUCKET_PATH_PREFIX}/{point_id}.pdf"
                preview_public_url = await StorageService.upload_bytes(
                    preview_object_name,
                    pdf_bytes,
                    content_type="application/pdf",
                )
                # Persist a stable preview pointer. Signed URLs are minted on demand.
                preview_url = preview_public_url
                preview_type = "pdf"
                preview_status = "complete"
            except HTTPException as exc:
                preview_status = "failed"
                preview_error = exc.detail if isinstance(exc.detail, str) else str(exc.detail)
            except Exception as exc:
                preview_status = "failed"
                preview_error = str(exc)

        target_collection = (
            settings.QDRANT_COLLECTION_TIER_1 if is_global_upload
            else settings.QDRANT_COLLECTION_TIER_2
        )
        ingestion_status = "uploaded"

        payload: Dict[str, Any] = {
            "record_type": "file",
            "filename": filename,
            "file_type": file_type,
            "type": file_type,
            "url": file_url,
            "is_global": is_global_upload,
            "tags": tags,
            "size": size_str,
            "size_bytes": file_size,
            "upload_date": upload_date,
            "uploaded_at": upload_date,
            "ai_enabled": False,
            "raw_content": "",
            "cleaned_content": "",
            "content": "",
            "content_preview": "",
            "ingestion_status": ingestion_status,
            "ingestion_error": None,
            "extracted_char_count": 0,
            "chunk_count": 0,
            "embedding_model": settings.EMBEDDING_MODEL,
            "embedding_tokens_estimate": 0,
            "embedding_cost_estimate_usd": 0.0,
            "chunk_profiles": {"micro": 0, "macro": 0},
            "bm25_enabled": False,
            "preview_url": preview_url,
            "preview_object_name": preview_object_name,
            "preview_type": preview_type,
            "preview_status": preview_status,
            "preview_error": preview_error,
            "ocr_enabled": False,
            "ocr_status": "not_requested",
            "ocr_text_present": False,
            "ocr_confidence": None,
            "multimodal_status": "pending",
            "multimodal_enabled": True,
            "ocr_processed": False,
            "vision_processed": False,
            "visual_chunk_count": 0,
            "vision_enabled": bool(settings.RILEY_VISION_ENABLED),
            "vision_status": "not_requested",
            "vision_segments_annotated": 0,
            "analysis_status": "not_requested",
            "analysis_completed_at": None,
            "analysis_error": None,
        }
        if not is_global_upload:
            payload["client_id"] = tenant_id
        if is_image:
            payload["content"] = "[Image file — Riley Memory is off by default; enable it to start OCR/ingestion]"
            payload["content_preview"] = payload["content"]

        placeholder_vector = [0.0] * int(settings.EMBEDDING_DIM)
        await vector_service.client.upsert(
            collection_name=target_collection,
            points=[PointStruct(id=point_id, vector=placeholder_vector, payload=payload)],
        )

        logger.info(
            "upload_staged file_id=%s tenant=%s file_type=%s collection=%s ai_enabled=%s status=%s",
            point_id,
            tenant_id,
            file_type,
            target_collection,
            False,
            ingestion_status,
        )
        return {
            "id": point_id,
            "url": file_url,
            "filename": filename,
            "type": file_type,
            "preview_status": preview_status,
            "preview_error": preview_error,
        }
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process upload: {exc}"
        ) from exc


async def reindex_existing_file(
    file_id: str,
    collection_name: str,
) -> Optional[str]:
    """Create and enqueue a durable reindex job for an existing file."""
    points = await vector_service.client.retrieve(
        collection_name=collection_name,
        ids=[file_id],
        with_payload=True,
        with_vectors=False,
    )
    if not points:
        raise HTTPException(status_code=404, detail=f"File with ID {file_id} not found")

    point = points[0]
    payload = point.payload or {}
    if payload.get("record_type") == "chunk":
        return None

    status_now = str(payload.get("ingestion_status") or "uploaded").strip().lower()
    if status_now in {"queued", "processing"}:
        return None

    file_type = (payload.get("file_type") or payload.get("type") or "").lower()
    if file_type not in TEXT_NATIVE_EXTENSIONS and file_type not in IMAGE_EXTENSIONS:
        await vector_service.client.set_payload(
            collection_name=collection_name,
            payload={
                "ingestion_status": "failed",
                "ingestion_error": f"Unsupported file type for multimodal ingestion: {file_type or 'unknown'}",
            },
            points=[file_id],
        )
        return None

    file_url = str(payload.get("url") or "").strip()
    if not file_url:
        raise HTTPException(status_code=400, detail=f"File {file_id} has no URL")

    file_bytes = await StorageService.download_file(file_url)
    source_content_sha256 = _content_sha256_hex(file_bytes)
    existing_indexed_sha = str(payload.get("indexed_content_sha256") or "").strip()
    existing_chunk_count = int(payload.get("chunk_count") or 0)

    campaign_id = str(payload.get("client_id") or "global")
    if status_now == "indexed" and existing_chunk_count > 0 and not existing_indexed_sha:
        await vector_service.client.set_payload(
            collection_name=collection_name,
            payload={"indexed_content_sha256": source_content_sha256},
            points=[file_id],
        )
        logger.info(
            "legacy_index_hash_backfilled file_id=%s campaign_id=%s collection=%s content_sha256=%s",
            file_id,
            campaign_id,
            collection_name,
            source_content_sha256,
        )
        logger.info(
            "reindex_skipped_already_indexed file_id=%s campaign_id=%s collection=%s content_sha256=%s",
            file_id,
            campaign_id,
            collection_name,
            source_content_sha256,
        )
        return None

    if (
        status_now == "indexed"
        and existing_chunk_count > 0
        and existing_indexed_sha
        and existing_indexed_sha == source_content_sha256
    ):
        logger.info(
            "reindex_skipped_already_indexed file_id=%s campaign_id=%s collection=%s content_sha256=%s",
            file_id,
            campaign_id,
            collection_name,
            source_content_sha256,
        )
        return None

    if existing_indexed_sha and existing_indexed_sha != source_content_sha256:
        logger.info(
            "replacing_index_for_updated_document file_id=%s campaign_id=%s collection=%s old_sha=%s new_sha=%s",
            file_id,
            campaign_id,
            collection_name,
            existing_indexed_sha,
            source_content_sha256,
        )

    tenant_id = payload.get("client_id") or "global"
    logger.info(
        "reindex_enqueued_hash_mismatch file_id=%s campaign_id=%s collection=%s old_sha=%s new_sha=%s",
        file_id,
        campaign_id,
        collection_name,
        existing_indexed_sha or "<missing>",
        source_content_sha256,
    )
    job_id = await enqueue_ingestion_job(
        collection_name=collection_name,
        file_id=str(file_id),
        tenant_id=tenant_id,
    )
    return job_id


async def run_ingestion_job(
    *,
    job_id: str,
    file_id: str,
    collection_name: str,
) -> None:
    """Worker entrypoint: load job context and execute ingestion pipeline."""
    try:
        points = await vector_service.client.retrieve(
            collection_name=collection_name,
            ids=[file_id],
            with_payload=True,
            with_vectors=False,
        )
        if not points:
            raise HTTPException(status_code=404, detail=f"File with ID {file_id} not found")
        payload = points[0].payload or {}
        if payload.get("record_type") == "chunk":
            return

        current_job_id = payload.get("ingestion_job_id")
        if current_job_id and current_job_id != job_id:
            # Stale task; a newer job exists for this file.
            return

        attempt_count = int(payload.get("ingestion_job_attempt_count") or 0) + 1

        # Riley Memory is the source of truth for whether ingestion is allowed.
        if not bool(payload.get("ai_enabled", False)):
            now_iso = datetime.now().isoformat()
            current_status = str(payload.get("ingestion_status") or "uploaded").strip().lower()
            await vector_service.client.set_payload(
                collection_name=collection_name,
                payload={
                    "ingestion_status": "uploaded" if current_status in {"queued", "processing", "uploaded"} else current_status,
                    "ingestion_job_status": "cancelled",
                    "ingestion_job_completed_at": now_iso,
                    "ingestion_job_error_message": "Skipped because Riley Memory is disabled.",
                },
                points=[file_id],
            )
            return

        if attempt_count > MAX_RETRYABLE_INGESTION_ATTEMPTS:
            prior_failure_reason = str(payload.get("ingestion_failure_reason") or "").strip() or None
            prior_error_message = (
                str(payload.get("ingestion_job_error_message") or payload.get("ingestion_error") or "").strip()
                or None
            )
            prior_summary = ""
            if prior_failure_reason or prior_error_message:
                prior_summary = (
                    f" Previous failure: reason={prior_failure_reason or 'unknown'}; "
                    f"detail={(prior_error_message or 'n/a')[:240]}"
                )
            error_message = (
                f"Ingestion retry threshold exceeded ({attempt_count}>{MAX_RETRYABLE_INGESTION_ATTEMPTS})."
                f"{prior_summary}"
            )
            await mark_ingestion_job_permanently_failed(
                collection_name=collection_name,
                file_id=str(file_id),
                job_id=job_id,
                reason="attempt_threshold_exceeded",
                error_message=error_message,
                attempt_count=attempt_count,
                prior_failure_reason=prior_failure_reason,
                prior_error_message=prior_error_message,
            )
            logger.error(
                "ingestion_permanent_failure reason=attempt_threshold_exceeded job_id=%s file_id=%s collection=%s attempt_count=%s threshold=%s prior_reason=%s",
                job_id,
                file_id,
                collection_name,
                attempt_count,
                MAX_RETRYABLE_INGESTION_ATTEMPTS,
                prior_failure_reason or "none",
            )
            raise IngestionPermanentFailure("attempt_threshold_exceeded")

        filename = payload.get("filename") or "unknown"
        file_type = (payload.get("file_type") or payload.get("type") or "").lower()
        file_url = payload.get("url")
        if not file_url:
            error_message = "File has no URL; ingestion cannot proceed."
            await mark_ingestion_job_permanently_failed(
                collection_name=collection_name,
                file_id=str(file_id),
                job_id=job_id,
                reason="missing_file_url",
                error_message=error_message,
                attempt_count=attempt_count,
            )
            logger.error(
                "ingestion_permanent_failure reason=missing_file_url job_id=%s file_id=%s collection=%s",
                job_id,
                file_id,
                collection_name,
            )
            raise IngestionPermanentFailure("missing_file_url")

        if file_type not in TEXT_NATIVE_EXTENSIONS and file_type not in IMAGE_EXTENSIONS:
            error_message = f"Unsupported file type: {file_type or 'unknown'}"
            await mark_ingestion_job_permanently_failed(
                collection_name=collection_name,
                file_id=str(file_id),
                job_id=job_id,
                reason="unsupported_file_type",
                error_message=error_message,
                attempt_count=attempt_count,
            )
            logger.error(
                "ingestion_permanent_failure reason=unsupported_file_type job_id=%s file_id=%s collection=%s file_type=%s",
                job_id,
                file_id,
                collection_name,
                file_type,
            )
            raise IngestionPermanentFailure("unsupported_file_type")

        try:
            file_bytes = await StorageService.download_file(file_url)
        except HTTPException as exc:
            if exc.status_code == 404:
                error_message = str(exc.detail or "Storage object is missing.")
                await mark_ingestion_job_permanently_failed(
                    collection_name=collection_name,
                    file_id=str(file_id),
                    job_id=job_id,
                    reason="gcs_file_missing",
                    error_message=error_message,
                    attempt_count=attempt_count,
                )
                logger.error(
                    "ingestion_permanent_failure reason=gcs_file_missing job_id=%s file_id=%s collection=%s",
                    job_id,
                    file_id,
                    collection_name,
                )
                raise IngestionPermanentFailure("gcs_file_missing") from exc
            raise
        source_content_sha256 = _content_sha256_hex(file_bytes)
        existing_indexed_sha = str(payload.get("indexed_content_sha256") or "").strip()
        existing_chunk_count = int(payload.get("chunk_count") or 0)
        status_now = str(payload.get("ingestion_status") or "uploaded").strip().lower()
        campaign_id = str(payload.get("client_id") or "global")
        if status_now == "indexed" and existing_chunk_count > 0 and not existing_indexed_sha:
            completed_at = datetime.now().isoformat()
            await vector_service.client.set_payload(
                collection_name=collection_name,
                payload={
                    "indexed_content_sha256": source_content_sha256,
                    "ingestion_job_status": "indexed",
                    "ingestion_job_completed_at": completed_at,
                    "ingestion_job_error_message": None,
                },
                points=[file_id],
            )
            logger.info(
                "legacy_index_hash_backfilled file_id=%s campaign_id=%s collection=%s content_sha256=%s",
                file_id,
                campaign_id,
                collection_name,
                source_content_sha256,
            )
            logger.info(
                "reindex_skipped_already_indexed file_id=%s campaign_id=%s collection=%s content_sha256=%s",
                file_id,
                campaign_id,
                collection_name,
                source_content_sha256,
            )
            return
        if (
            status_now == "indexed"
            and existing_chunk_count > 0
            and existing_indexed_sha
            and existing_indexed_sha == source_content_sha256
        ):
            completed_at = datetime.now().isoformat()
            await vector_service.client.set_payload(
                collection_name=collection_name,
                payload={
                    "ingestion_job_status": "indexed",
                    "ingestion_job_completed_at": completed_at,
                    "ingestion_job_error_message": None,
                },
                points=[file_id],
            )
            logger.info(
                "reindex_skipped_already_indexed file_id=%s campaign_id=%s collection=%s content_sha256=%s",
                file_id,
                campaign_id,
                collection_name,
                source_content_sha256,
            )
            return
        tenant_id = payload.get("client_id") or "global"
        is_global = bool(payload.get("is_global")) or collection_name == get_settings().QDRANT_COLLECTION_TIER_1
        tags = payload.get("tags", []) if isinstance(payload.get("tags"), list) else []
        upload_date = payload.get("upload_date") or datetime.now().isoformat()
        size_bytes = payload.get("size_bytes") or len(file_bytes)
        size_str = payload.get("size") or _format_file_size(int(size_bytes))

        await _run_ingestion_pipeline(
            job_id=job_id,
            collection_name=collection_name,
            point_id=str(file_id),
            file_bytes=file_bytes,
            filename=filename,
            file_type=file_type,
            tenant_id=tenant_id,
            is_global=is_global,
            tags=tags,
            file_url=file_url,
            size_str=size_str,
            file_size=int(size_bytes),
            upload_date=upload_date,
            preview_url=payload.get("preview_url"),
            preview_type=payload.get("preview_type"),
            preview_status=payload.get("preview_status", "not_requested"),
            preview_error=payload.get("preview_error"),
            initial_ai_enabled=bool(payload.get("ai_enabled", False)),
            source_content_sha256=source_content_sha256,
            attempt_count=attempt_count,
        )
    except IngestionPermanentFailure:
        raise
    except Exception as exc:
        logger.warning(
            "ingestion_retryable_failure job_id=%s file_id=%s collection=%s error_type=%s error=%s",
            job_id,
            file_id,
            collection_name,
            type(exc).__name__,
            exc,
        )
        try:
            await _mark_job_failed(
                collection_name=collection_name,
                file_id=str(file_id),
                job_id=job_id,
                error_message=str(exc),
            )
        except Exception as status_exc:
            logger.exception(
                "ingestion_worker_status_update_failed job_id=%s file_id=%s error=%s",
                job_id,
                file_id,
                status_exc,
            )
        raise


def _is_safe_reenqueue_payload(payload: Dict[str, Any]) -> bool:
    if str(payload.get("record_type") or "").strip().lower() == "chunk":
        return False
    status_now = str(payload.get("ingestion_status") or "").strip().lower()
    job_status_now = str(payload.get("ingestion_job_status") or "").strip().lower()
    if status_now == "permanently_failed" or job_status_now == "permanently_failed":
        return False
    if not bool(payload.get("ai_enabled", False)):
        return False
    file_type = str(payload.get("file_type") or payload.get("type") or "").strip().lower()
    if file_type not in TEXT_NATIVE_EXTENSIONS and file_type not in IMAGE_EXTENSIONS:
        return False
    file_url = str(payload.get("url") or "").strip()
    if not file_url:
        return False
    return True


async def reset_stuck_processing_ingestion_jobs(
    *,
    older_than_minutes: int = 30,
    reenqueue_safe: bool = False,
    limit_per_collection: int = 500,
) -> Dict[str, Any]:
    """Reset stale ingestion jobs stuck in processing status.

    This utility is manual/admin-invoked and intentionally not auto-scheduled.
    """
    settings = get_settings()
    cutoff = datetime.now(timezone.utc) - timedelta(minutes=max(1, int(older_than_minutes)))
    collections = [
        settings.QDRANT_COLLECTION_TIER_2,
        settings.QDRANT_COLLECTION_TIER_1,
    ]
    summary: Dict[str, Any] = {
        "older_than_minutes": max(1, int(older_than_minutes)),
        "reenqueue_safe": bool(reenqueue_safe),
        "detected": 0,
        "reset_to_queued": 0,
        "reset_to_failed": 0,
        "reenqueued": 0,
        "errors": 0,
    }

    processing_filter = Filter(
        must=[
            FieldCondition(key="record_type", match=MatchValue(value="file")),
            FieldCondition(key="ingestion_status", match=MatchValue(value="processing")),
        ]
    )

    for collection_name in collections:
        offset = None
        while True:
            points, next_offset = await vector_service.client.scroll(
                collection_name=collection_name,
                scroll_filter=processing_filter,
                limit=max(1, int(limit_per_collection)),
                offset=offset,
                with_payload=True,
                with_vectors=False,
            )
            if not points:
                break

            for point in points:
                point_id = str(point.id)
                payload = point.payload or {}
                started_at = _parse_payload_datetime(
                    payload.get("ingestion_started_at") or payload.get("ingestion_job_started_at")
                )
                if started_at is None or started_at > cutoff:
                    continue

                summary["detected"] += 1
                logger.warning(
                    "ingestion_stuck_processing_detected file_id=%s collection=%s ingestion_started_at=%s cutoff=%s",
                    point_id,
                    collection_name,
                    payload.get("ingestion_started_at") or payload.get("ingestion_job_started_at"),
                    cutoff.isoformat(),
                )
                try:
                    if _is_safe_reenqueue_payload(payload):
                        reset_payload = {
                            "ingestion_status": "queued",
                            "ingestion_error": "Reset from stale processing state; awaiting retry.",
                            "ingestion_job_status": "queued",
                            "ingestion_started_at": None,
                            "ingestion_job_started_at": None,
                        }
                        await vector_service.client.set_payload(
                            collection_name=collection_name,
                            payload=reset_payload,
                            points=[point_id],
                        )
                        summary["reset_to_queued"] += 1
                        logger.info(
                            "ingestion_status_reset action=reset_to_queued file_id=%s collection=%s",
                            point_id,
                            collection_name,
                        )
                        if reenqueue_safe:
                            tenant_id = str(payload.get("client_id") or "global")
                            try:
                                await enqueue_ingestion_job(
                                    collection_name=collection_name,
                                    file_id=point_id,
                                    tenant_id=tenant_id,
                                )
                                summary["reenqueued"] += 1
                                logger.info(
                                    "ingestion_status_reset action=reenqueued file_id=%s collection=%s",
                                    point_id,
                                    collection_name,
                                )
                            except Exception as exc:
                                summary["errors"] += 1
                                logger.exception(
                                    "ingestion_status_reset action=reenqueue_failed file_id=%s collection=%s error=%s",
                                    point_id,
                                    collection_name,
                                    exc,
                                )
                    else:
                        await mark_ingestion_job_permanently_failed(
                            collection_name=collection_name,
                            file_id=point_id,
                            job_id=str(payload.get("ingestion_job_id") or "").strip() or None,
                            reason="stuck_processing_unrecoverable",
                            error_message="Stale processing state reset; payload unsafe to re-enqueue.",
                            attempt_count=int(payload.get("ingestion_job_attempt_count") or 0),
                        )
                        summary["reset_to_failed"] += 1
                        logger.info(
                            "ingestion_status_reset action=reset_to_failed file_id=%s collection=%s",
                            point_id,
                            collection_name,
                        )
                except Exception as exc:
                    summary["errors"] += 1
                    logger.exception(
                        "ingestion_status_reset action=error file_id=%s collection=%s error=%s",
                        point_id,
                        collection_name,
                        exc,
                    )

            if next_offset is None:
                break
            offset = next_offset

    return summary


async def backfill_reindex_supported_files(
    tenant_id: Optional[str] = None,
    include_global: bool = True,
    limit: int = 500,
) -> Dict[str, int]:
    """Backfill existing file points by enqueueing durable Phase 1A jobs."""
    settings = get_settings()
    totals = {"seen": 0, "enqueued": 0, "skipped": 0, "failed": 0}

    targets: List[Tuple[str, Optional[Filter]]] = []
    tenant_filter = None
    if tenant_id and tenant_id != "global":
        tenant_filter = Filter(
            must=[FieldCondition(key="client_id", match=MatchValue(value=tenant_id))]
        )
    targets.append((settings.QDRANT_COLLECTION_TIER_2, tenant_filter))
    if include_global:
        global_filter = Filter(must=[FieldCondition(key="is_global", match=MatchValue(value=True))])
        targets.append((settings.QDRANT_COLLECTION_TIER_1, global_filter))

    for collection_name, collection_filter in targets:
        points, _ = await vector_service.client.scroll(
            collection_name=collection_name,
            scroll_filter=collection_filter,
            limit=limit,
            with_payload=True,
            with_vectors=False,
        )
        for point in points:
            totals["seen"] += 1
            payload = point.payload or {}
            if payload.get("record_type") == "chunk":
                totals["skipped"] += 1
                continue
            file_type = (payload.get("file_type") or payload.get("type") or "").lower()
            if file_type not in TEXT_NATIVE_EXTENSIONS:
                totals["skipped"] += 1
                continue
            try:
                job_id = await reindex_existing_file(str(point.id), collection_name)
                if job_id:
                    totals["enqueued"] += 1
                else:
                    totals["skipped"] += 1
            except Exception:
                totals["failed"] += 1
                logger.exception("backfill_reindex_failed file_id=%s collection=%s", point.id, collection_name)
    return totals


async def backfill_indexed_content_hashes(
    *,
    tenant_id: Optional[str] = None,
    include_global: bool = True,
    batch_size: int = 100,
    max_points_per_collection: Optional[int] = None,
) -> Dict[str, int]:
    """
    One-time metadata backfill for indexed files missing indexed_content_sha256.

    Safety guarantees:
    - metadata-only payload update on parent file points
    - no ingestion enqueue
    - no embedding generation
    - no chunk point mutation
    """
    settings = get_settings()
    effective_batch_size = max(50, min(100, int(batch_size or 100)))
    totals = {
        "scanned": 0,
        "processed": 0,
        "skipped": 0,
        "errors": 0,
    }

    targets: List[Tuple[str, Optional[Filter]]] = []
    tenant_filter = None
    if tenant_id and tenant_id != "global":
        tenant_filter = Filter(
            must=[FieldCondition(key="client_id", match=MatchValue(value=tenant_id))],
            must_not=[FieldCondition(key="record_type", match=MatchValue(value="chunk"))],
        )
    else:
        tenant_filter = Filter(
            must_not=[FieldCondition(key="record_type", match=MatchValue(value="chunk"))],
        )
    targets.append((settings.QDRANT_COLLECTION_TIER_2, tenant_filter))
    if include_global:
        global_filter = Filter(
            must=[FieldCondition(key="is_global", match=MatchValue(value=True))],
            must_not=[FieldCondition(key="record_type", match=MatchValue(value="chunk"))],
        )
        targets.append((settings.QDRANT_COLLECTION_TIER_1, global_filter))

    logger.info(
        "backfill_started tenant_id=%s include_global=%s batch_size=%s max_points_per_collection=%s",
        tenant_id,
        include_global,
        effective_batch_size,
        max_points_per_collection,
    )
    max_file_size_bytes = int(settings.MAX_UPLOAD_MB) * 1024 * 1024

    def _resolve_file_size_bytes_from_gcs(file_url: str) -> int:
        client = StorageService._get_client()
        bucket_name, _, decoded_blob_name = StorageService._parse_gcs_location(file_url)
        bucket = client.bucket(bucket_name)
        blob = bucket.blob(decoded_blob_name)
        blob.reload()
        return int(blob.size or 0)

    for collection_name, base_filter in targets:
        offset: Any = None
        scanned_for_collection = 0
        batch_num = 0
        while True:
            if max_points_per_collection and scanned_for_collection >= max_points_per_collection:
                break

            try:
                points, next_offset = await vector_service.client.scroll(
                    collection_name=collection_name,
                    scroll_filter=base_filter,
                    limit=effective_batch_size,
                    offset=offset,
                    with_payload=True,
                    with_vectors=False,
                )
            except Exception as exc:
                message = str(exc).lower()
                if "index required but not found" not in message:
                    raise
                # Legacy fallback when record_type payload index is missing.
                legacy_filter = Filter(
                    must=(base_filter.must or []) if base_filter else [],
                    should=(base_filter.should or []) if base_filter else [],
                )
                points, next_offset = await vector_service.client.scroll(
                    collection_name=collection_name,
                    scroll_filter=legacy_filter,
                    limit=effective_batch_size,
                    offset=offset,
                    with_payload=True,
                    with_vectors=False,
                )
                points = [point for point in points if (point.payload or {}).get("record_type") != "chunk"]

            if not points:
                break

            batch_num += 1
            batch_processed = 0
            batch_skipped = 0
            batch_errors = 0
            scanned_for_collection += len(points)
            totals["scanned"] += len(points)

            for point in points:
                payload = point.payload or {}
                if payload.get("record_type") == "chunk":
                    totals["skipped"] += 1
                    batch_skipped += 1
                    continue

                status_now = str(payload.get("ingestion_status") or "").strip().lower()
                chunk_count = int(payload.get("chunk_count") or 0)
                existing_hash = str(payload.get("indexed_content_sha256") or "").strip()
                if status_now != "indexed" or chunk_count <= 0 or existing_hash:
                    totals["skipped"] += 1
                    batch_skipped += 1
                    continue

                file_id = str(point.id)
                file_url = str(payload.get("url") or "").strip()
                campaign_id = str(payload.get("client_id") or ("global" if bool(payload.get("is_global")) else "unknown"))
                if not file_url:
                    totals["errors"] += 1
                    batch_errors += 1
                    logger.warning(
                        "backfill_errors file_id=%s campaign_id=%s collection=%s reason=missing_file_url",
                        file_id,
                        campaign_id,
                        collection_name,
                    )
                    continue

                raw_size = payload.get("size_bytes")
                file_size_bytes = 0
                if isinstance(raw_size, (int, float)):
                    file_size_bytes = int(raw_size)
                elif isinstance(raw_size, str):
                    try:
                        file_size_bytes = int(float(raw_size.strip()))
                    except Exception:
                        file_size_bytes = 0

                if file_size_bytes <= 0:
                    try:
                        file_size_bytes = await run_in_threadpool(_resolve_file_size_bytes_from_gcs, file_url)
                    except Exception as exc:
                        totals["errors"] += 1
                        batch_errors += 1
                        logger.warning(
                            "backfill_errors file_id=%s campaign_id=%s collection=%s error=%s",
                            file_id,
                            campaign_id,
                            collection_name,
                            type(exc).__name__,
                        )
                        continue

                if file_size_bytes > max_file_size_bytes:
                    totals["skipped"] += 1
                    batch_skipped += 1
                    logger.info(
                        "backfill_skipped_large_file file_id=%s campaign_id=%s collection=%s size=%s threshold=%s",
                        file_id,
                        campaign_id,
                        collection_name,
                        file_size_bytes,
                        max_file_size_bytes,
                    )
                    continue

                try:
                    file_bytes = await StorageService.download_file(file_url)
                    source_content_sha256 = _content_sha256_hex(file_bytes)
                    await vector_service.client.set_payload(
                        collection_name=collection_name,
                        payload={"indexed_content_sha256": source_content_sha256},
                        points=[file_id],
                    )
                    totals["processed"] += 1
                    batch_processed += 1
                except Exception as exc:
                    totals["errors"] += 1
                    batch_errors += 1
                    logger.warning(
                        "backfill_errors file_id=%s campaign_id=%s collection=%s error=%s",
                        file_id,
                        campaign_id,
                        collection_name,
                        type(exc).__name__,
                    )

            logger.info(
                "backfill_processed_count collection=%s batch=%s batch_processed=%s total_processed=%s",
                collection_name,
                batch_num,
                batch_processed,
                totals["processed"],
            )
            logger.info(
                "backfill_skipped_count collection=%s batch=%s batch_skipped=%s total_skipped=%s",
                collection_name,
                batch_num,
                batch_skipped,
                totals["skipped"],
            )
            if batch_errors:
                logger.info(
                    "backfill_errors collection=%s batch=%s batch_errors=%s total_errors=%s",
                    collection_name,
                    batch_num,
                    batch_errors,
                    totals["errors"],
                )

            offset = next_offset
            if offset is None:
                break

    logger.info(
        "backfill_completed scanned=%s processed=%s skipped=%s errors=%s",
        totals["scanned"],
        totals["processed"],
        totals["skipped"],
        totals["errors"],
    )
    return totals


async def delete_file(file_id: str, collection_name: str) -> None:
    """
    Permanently delete a file from GCS and Qdrant.
    
    Args:
        file_id: The Qdrant point ID (UUID) of the file to delete
        collection_name: Name of the Qdrant collection
        
    Raises:
        HTTPException: If file not found or deletion fails
    """
    # Retrieve the point to get file metadata
    try:
        points = await vector_service.client.retrieve(
            collection_name=collection_name,
            ids=[file_id],
            with_payload=True,
            with_vectors=False
        )
        
        if not points:
            raise HTTPException(
                status_code=404,
                detail=f"File with ID {file_id} not found in Qdrant"
            )
        
        point = points[0]
        payload = point.payload or {}
        file_url = payload.get("url")
        
        if not file_url:
            raise HTTPException(
                status_code=404,
                detail=f"File URL not found for ID {file_id}"
            )
        
        # Delete file from GCS
        try:
            await StorageService.delete_file(file_url)
        except HTTPException as exc:
            # If file not found in GCS, log but continue with Qdrant deletion
            # (file might have been manually deleted from GCS)
            if exc.status_code == 404:
                pass  # Continue with Qdrant deletion
            else:
                raise
        
        # Delete point from Qdrant
        await vector_service.client.delete(
            collection_name=collection_name,
            points_selector=[file_id]
        )

        # Delete chunk points linked to this parent file
        await _delete_chunk_points(collection_name, file_id)
        
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to delete file: {exc}"
        ) from exc


async def untag_file(file_id: str, tag: str | None, collection_name: str) -> None:
    """
    Remove a specific tag from a file's metadata in Qdrant, or clear all tags.
    
    Args:
        file_id: The Qdrant point ID (UUID) of the file
        tag: The tag to remove (e.g., "Messaging", "Research"). If None, clears all tags.
        collection_name: Name of the Qdrant collection
        
    Raises:
        HTTPException: If file not found or update fails
    """
    # Retrieve the point to get current payload
    try:
        points = await vector_service.client.retrieve(
            collection_name=collection_name,
            ids=[file_id],
            with_payload=True,
            with_vectors=False
        )
        
        if not points:
            raise HTTPException(
                status_code=404,
                detail=f"File with ID {file_id} not found in Qdrant"
            )
        
        point = points[0]
        payload = point.payload or {}
        current_tags = payload.get("tags", [])
        
        # If tag is None, clear all tags
        if tag is None:
            updated_tags = []
        else:
            # Remove the specific tag if it exists
            if isinstance(current_tags, list) and tag in current_tags:
                updated_tags = [t for t in current_tags if t != tag]
            else:
                # Tag doesn't exist, but that's okay - just return success
                return
        
        # Update the payload with the new tags list
        await vector_service.client.set_payload(
            collection_name=collection_name,
            payload={"tags": updated_tags},
            points=[file_id]
        )
        
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to untag file: {exc}"
        ) from exc

