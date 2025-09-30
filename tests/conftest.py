from pathlib import Path

import pytest


def _make_docx(path: Path, text: str) -> None:
    import docx

    d = docx.Document()
    for line in text.split("\n"):
        d.add_paragraph(line)
    d.save(str(path))


def _make_pptx(path: Path, slides: list[str]) -> None:
    from pptx import Presentation

    prs = Presentation()
    for stext in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = (stext.split("\n", 1)[0] or "Slide")[:100]
        slide.placeholders[1].text = stext
    prs.save(str(path))


def _make_pdf(path: Path, text: str) -> None:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=letter)
    width, height = letter
    y = height - 72
    for line in text.split("\n"):
        c.drawString(72, y, line[:1000])
        y -= 14
        if y < 72:
            c.showPage()
            y = height - 72
    c.save()


@pytest.fixture(scope="session")
def sample_campaign(tmp_path_factory):
    root = tmp_path_factory.mktemp("sample_campaign")
    _make_docx(
        root / "fundraising_memo.docx",
        "Urgent: Small-dollar fundraising push.\nDonate at donate@example.org\nCall +1 202 555 0134",
    )
    _make_pptx(
        root / "press_announcement.pptx",
        [
            "Press Statement",
            "We will announce the plan today. Media briefing scheduled.",
        ],
    )
    _make_pdf(
        root / "policy_brief.pdf",
        "Policy: Our healthcare plan reduces costs. Contact press@campaign.org for details.",
    )
    (root / "field_volunteer.txt").write_text(
        "Volunteer to canvass and phonebank this weekend!\nMeet at 123 Main Street.",
        encoding="utf-8",
    )
    return root


@pytest.fixture
def ci_out(tmp_path):
    out = tmp_path / ".ci_out"
    out.mkdir()
    return out
