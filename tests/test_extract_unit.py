from __future__ import annotations
from pathlib import Path
import pandas as pd
from pipeline.extract import (
    read_docx,
    read_pdf,
    read_pptx,
    summarize_xlsx,
    sniff_and_read,
)


def _make_docx(path: Path, text: str) -> None:
    import docx

    d = docx.Document()
    for line in text.split("\n"):
        d.add_paragraph(line)
    d.save(str(path))


def _make_pptx(path: Path, slides: list[str]) -> None:
    from pptx import Presentation

    prs = Presentation()
    for stext in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = (stext.split("\n", 1)[0] or "Slide")[:100]
        slide.placeholders[1].text = stext
    prs.save(str(path))


def _make_pdf(path: Path, text: str) -> None:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=letter)
    width, height = letter
    y = height - 72
    for line in text.split("\n"):
        c.drawString(72, y, line[:1000])
        y -= 14
        if y < 72:
            c.showPage()
            y = height - 72
    c.save()


def _make_xlsx(path: Path) -> None:
    # 2 sheets with columns that resemble media list/editorial calendar
    with pd.ExcelWriter(path) as w:
        df1 = pd.DataFrame(
            {
                "outlet": ["XYZ News"],
                "reporter": ["Alice"],
                "email": ["a@example.com"],
                "beat": ["health"],
            }
        )
        df1.to_excel(w, index=False, sheet_name="Media")

        df2 = pd.DataFrame(
            {
                "month": ["January", "February"],
                "theme": ["Launch", "Follow-up"],
                "status": ["draft", "final"],
            }
        )
        df2.to_excel(w, index=False, sheet_name="Editorial Calendar")


def test_read_docx(tmp_path: Path):
    p = tmp_path / "a.docx"
    _make_docx(p, "Hello world\nSecond line")
    out = read_docx(p)
    assert "Hello world" in out
    assert "Second line" in out


def test_read_pdf(tmp_path: Path):
    p = tmp_path / "a.pdf"
    _make_pdf(p, "Policy brief\nContact press@org.org")
    out = read_pdf(p)
    assert "Policy brief" in out
    assert "press@org.org" in out


def test_read_pptx(tmp_path: Path):
    p = tmp_path / "a.pptx"
    _make_pptx(p, ["Slide Title\nBody text", "Another Title\nMore body"])
    slides = read_pptx(p)
    assert isinstance(slides, dict)
    assert 1 in slides
    assert "title" in slides[1] and "body" in slides[1]
    assert "Slide Title" in slides[1]["title"]


def test_summarize_xlsx(tmp_path: Path):
    p = tmp_path / "a.xlsx"
    _make_xlsx(p)
    schema = summarize_xlsx(p)
    assert "sheets" in schema and isinstance(schema["sheets"], list)
    assert any(s["sheet"].lower().startswith("media") for s in schema["sheets"])


def test_sniff_and_read_all(tmp_path: Path):
    dp = tmp_path / "docx.docx"
    _make_docx(dp, "DocX here")
    pp = tmp_path / "slides.pptx"
    _make_pptx(pp, ["One\nTwo"])
    pf = tmp_path / "a.pdf"
    _make_pdf(pf, "PDF text")
    tx = tmp_path / "plain.txt"
    tx.write_text("Hello TXT", encoding="utf-8")
    xx = tmp_path / "schema.xlsx"
    _make_xlsx(xx)

    assert "DocX here" in sniff_and_read(dp)
    assert "One" in sniff_and_read(pp)
    assert "PDF text" in sniff_and_read(pf)
    assert "Hello TXT" in sniff_and_read(tx)
    # xlsx returns empty string by design in sniff; summarize_xlsx used in pipeline.extract main
    out = sniff_and_read(xx)
    assert isinstance(out, str)
    assert "sheets" in out  # since summarize_xlsx is returned as JSON
