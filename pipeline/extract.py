import sys
import json
import logging
from pathlib import Path
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation

# Suppress noisy PDF logging
for name in ("pdfminer", "pdfplumber"):
    logging.getLogger(name).setLevel(logging.ERROR)


# --- Safe readers -------------------------------------------------------


def read_docx(p: Path) -> str:
    """Extract text from a DOCX file safely."""
    try:
        d = Document(p.as_posix())
        return "\n".join(para.text for para in d.paragraphs).strip()
    except Exception as e:
        return f"[DOCX_READ_ERROR: {e}]"


def read_pdf(p: Path) -> str:
    """Extract text from a PDF using pdfplumber, page by page."""
    out = []
    try:
        with pdfplumber.open(p.as_posix()) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                out.append(text)
        return "\n\n".join(out).strip()
    except Exception as e:
        return f"[PDF_READ_ERROR: {e}]"


def read_pptx(p: Path) -> dict:
    """Extract slide titles and body text from a PPTX file."""
    slides = {}
    try:
        prs = Presentation(p.as_posix())
        for i, s in enumerate(prs.slides, start=1):
            title = ""
            body = []
            for shp in s.shapes:
                if hasattr(shp, "text") and shp.text:
                    if not title:
                        title = shp.text.strip()
                    else:
                        body.append(shp.text.strip())
            slides[i] = {"title": title, "body": "\n".join(body)}
    except Exception as e:
        slides = {"error": str(e)}
    return slides


def summarize_xlsx(p: Path) -> dict:
    """Summarize an Excel workbook: sheet names, row/col counts, and top columns."""
    out = {"sheets": [], "pii_any": False}
    try:
        xls = pd.ExcelFile(p.as_posix())
        for name in xls.sheet_names:
            try:
                df = xls.parse(name)
                cols = [str(c) for c in df.columns]
                out["sheets"].append(
                    {
                        "sheet": name,
                        "rows": int(df.shape[0]),
                        "cols": int(df.shape[1]),
                        "columns": cols[:25] + (["..."] if len(cols) > 25 else []),
                    }
                )
            except Exception as e:
                out["sheets"].append({"sheet": name, "error": str(e)})
    except Exception as e:
        out["sheets"].append({"error": f"Failed to open workbook: {e}"})
    return out


# --- Generic entry point ------------------------------------------------


def sniff_and_read(path: Path) -> str:
    """
    Generic extractor for CI/CD and smoke tests.
    Returns plain text from PDFs, DOCX, PPTX, TXT, or schema string from XLSX.
    """
    ext = path.suffix.lower()
    if ext == ".pdf":
        return read_pdf(path)
    if ext == ".docx":
        return read_docx(path)
    if ext in (".pptx", ".ppt"):
        slides = read_pptx(path)
        if "error" in slides:
            return f"[PPTX_READ_ERROR: {slides['error']}]"
        return "\n".join(
            [f"{v.get('title','')} {v.get('body','')}" for v in slides.values()]
        )
    if ext in (".xlsx", ".xls"):
        schema = summarize_xlsx(path)
        return json.dumps(schema, ensure_ascii=False)
    if ext in (".txt", ".md"):
        try:
            return path.read_text(encoding="utf-8", errors="ignore").strip()
        except Exception:
            return path.read_text(errors="ignore").strip()
    return ""


# --- Campaign-level driver ----------------------------------------------


def main(campaign_root: str):
    """
    Campaign-level extractor:
    Reads inventory CSV and writes extracted.jsonl
    with structured content per file.
    """
    camp = Path(campaign_root).name
    inv = Path("data/processed") / f"inventory_{camp}.csv"

    # --- Guard: inventory existence and readability ---
    if not inv.exists():
        print(f"[extract] skip {camp} (inventory not found: {inv})")
        return
    if inv.stat().st_size == 0:
        print(f"[extract] skip {camp} (0 files in inventory)")
        return

    try:
        df = pd.read_csv(inv)
    except pd.errors.EmptyDataError:
        print(f"[extract] skip {camp} (inventory empty/unreadable: {inv})")
        return
    except Exception as e:
        print(f"[extract] skip {camp} (failed to read inventory: {e})")
        return

    if df.empty:
        print(f"[extract] skip {camp} (inventory has no rows)")
        return

    out_dir = Path("data/processed") / camp
    out_dir.mkdir(parents=True, exist_ok=True)
    out_jsonl = out_dir / "extracted.jsonl"

    with out_jsonl.open("w", encoding="utf-8") as f:
        for _, row in df.iterrows():
            p = Path(row["source_path"])
            rec = {
                "source_path": row.get("source_path", ""),
                "sha256": row.get("sha256", ""),
                "ext": row.get("ext", ""),
            }
            try:
                if p.suffix.lower() == ".docx":
                    rec["extract"] = {"type": "text", "text": read_docx(p)}
                elif p.suffix.lower() == ".pdf":
                    rec["extract"] = {"type": "text", "text": read_pdf(p)}
                elif p.suffix.lower() in (".pptx", ".ppt"):
                    rec["extract"] = {"type": "slides", "slides": read_pptx(p)}
                elif p.suffix.lower() in (".xlsx", ".xls"):
                    rec["extract"] = {
                        "type": "xlsx_schema",
                        "schema": summarize_xlsx(p),
                    }
                elif p.suffix.lower() in (".txt", ".md"):
                    rec["extract"] = {"type": "text", "text": sniff_and_read(p)}
                else:
                    rec["extract"] = {"type": "skip"}
            except Exception as e:
                rec["error"] = str(e)
            f.write(json.dumps(rec, ensure_ascii=False) + "\n")

    print(f"[extract] wrote {out_jsonl}")


if __name__ == "__main__":
    main(sys.argv[1])
